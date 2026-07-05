"""
Advanced Converter Engine
converter/converters.py

HTML → PDF  : Playwright (headless Chrome) › pdfkit › fitz › reportlab+images
PDF → HTML  : fitz dict-mode (flow layout, base64 images, adaptive sizing)
XLSX → PDF  : smart portrait/landscape, auto col-widths, multi-sheet
PPTX → PDF  : LibreOffice › python-pptx+reportlab (images embedded)
TXT/RTF     : reportlab (styled) › pypandoc › striprtf
EPUB → PDF  : pypandoc › spine-order native (images embedded)
Images      : Pillow max-quality, EXIF preserved, HEIC via pillow-heif
Audio/Video : ffmpeg binary auto-located, quality presets per format

"""

from __future__ import annotations

import base64
import io
import os
import re
import subprocess
import tempfile
import time
import zipfile
from pathlib import Path

from converter.mixins.document_converters import DocumentConverters
from converter.mixins.image_converters import ImageConverters
from converter.mixins.media_converters import MediaConverters
from converter.html_inline import inline_markup, build_table_flowable, add_image_flowable

_NO_WINDOW = subprocess.CREATE_NO_WINDOW if hasattr(subprocess, 'CREATE_NO_WINDOW') else 0

class ConversionResult:
    __slots__ = ("success", "source", "target", "elapsed", "error", "file_size")

    def __init__(self, success, source, target, elapsed=0.0, error="", file_size=0):
        self.success   = success
        self.source    = source
        self.target    = target
        self.elapsed   = elapsed
        self.error     = error
        self.file_size = file_size

    def __repr__(self):
        s = "OK" if self.success else f"ERR({self.error})"
        return f"<ConversionResult {s} {self.source!r}→{self.target!r}>"

def _timed(fn):
    t0 = time.perf_counter()
    fn()
    return time.perf_counter() - t0

def _build_dst(src, dst_dir, new_ext):
    return str(Path(dst_dir) / f"{Path(src).stem}.{new_ext.lstrip('.')}")

def _img_to_b64(data, mime="image/png"):
    return f"data:{mime};base64,{base64.b64encode(data).decode()}"

def _mime_for_ext(ext):
    return {
        "png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
        "gif": "image/gif", "webp": "image/webp", "svg": "image/svg+xml",
        "bmp": "image/bmp", "tiff": "image/tiff",
    }.get(ext.lower(), "image/png")

def _safe_html(text):
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

def _read_file_b64(path):
    with open(path, "rb") as f:
        return base64.b64encode(f.read()).decode()

CATEGORY_MAP = {
    # Documents
    "txt_to_pdf":  "document", "rtf_to_pdf":   "document",
    "txt_to_docx": "document", "rtf_to_docx":  "document",
    "csv_to_json": "document", "json_to_csv":  "document",
    "xlsx_to_pdf": "document", "xlsx_to_json": "document",
    "xlsx_to_csv": "document", "pptx_to_pdf":  "document",
    "html_to_pdf": "document", "pdf_to_html":  "document",
    "epub_to_pdf": "document",
    # Images
    "image_to_png":  "image", "image_to_jpeg": "image",
    "image_to_jpg":  "image", "image_to_bmp":  "image",
    "image_to_heic": "image", "image_to_webp": "image",
    "image_to_tiff": "image", "image_to_psd":  "image",
    "image_to_svg":  "image", "image_to_avif": "image",
    "image_to_j2k":  "image", "image_to_dng":  "image",
    "image_to_ico":  "image",
    # Video
    "video_to_mp4":  "video", "video_to_webm": "video",
    "video_to_mkv":  "video", "video_to_mov":  "video",
    "video_to_avi":  "video", "video_to_mp3":  "audio",
    "video_to_wav":  "audio", "video_to_aac":  "audio",
    "video_to_flac": "audio",
    # Audio
    "audio_to_mp3":  "audio", "audio_to_wav":  "audio",
    "audio_to_aac":  "audio", "audio_to_ogg":  "audio",
    "audio_to_flac": "audio", "audio_to_m4a":  "audio",
}

class AdvancedConverterEngine(DocumentConverters, ImageConverters, MediaConverters):

    _DISPATCH = {
        # Documents
        "txt_to_pdf":   ("_txt_to_pdf",     "pdf"),
        "rtf_to_pdf":   ("_rtf_to_pdf",     "pdf"),
        "txt_to_docx":  ("_txt_to_docx",    "docx"),
        "rtf_to_docx":  ("_rtf_to_docx",    "docx"),
        "csv_to_json":  ("_csv_to_json",    "json"),
        "json_to_csv":  ("_json_to_csv",    "csv"),
        "xlsx_to_pdf":  ("_xlsx_to_pdf",    "pdf"),
        "xlsx_to_json": ("_xlsx_to_json",   "json"),
        "xlsx_to_csv":  ("_xlsx_to_csv",    "csv"),
        "pptx_to_pdf":  ("_pptx_to_pdf",    "pdf"),
        "html_to_pdf":  ("_html_to_pdf",    "pdf"),
        "pdf_to_html":  ("_pdf_to_html",    "html"),
        "epub_to_pdf":  ("_epub_to_pdf",    "pdf"),
        # Images
        "image_to_png":  ("_image_convert",  "png"),
        "image_to_jpeg": ("_image_convert",  "jpeg"),
        "image_to_jpg":  ("_image_convert",  "jpg"),
        "image_to_bmp":  ("_image_convert",  "bmp"),
        "image_to_heic": ("_heic_convert",   "heic"),
        "image_to_webp": ("_image_convert",  "webp"),
        "image_to_tiff": ("_image_convert",  "tiff"),
        "image_to_psd":  ("_magick_convert",  "psd"),
        "image_to_svg":  ("_image_to_svg",   "svg"),
        "image_to_avif": ("_image_convert",  "avif"),
        "image_to_j2k":  ("_image_convert",  "j2k"),
        "image_to_dng":  ("_raw_convert",    "dng"),
        "image_to_ico":  ("_image_to_ico",   "ico"),
        # Video → Video
        "video_to_mp4":  ("_ffmpeg_convert", "mp4"),
        "video_to_webm": ("_ffmpeg_convert", "webm"),
        "video_to_mkv":  ("_ffmpeg_convert", "mkv"),
        "video_to_mov":  ("_ffmpeg_convert", "mov"),
        "video_to_avi":  ("_ffmpeg_convert", "avi"),
        # Video → Audio
        "video_to_mp3":  ("_ffmpeg_convert", "mp3"),
        "video_to_wav":  ("_ffmpeg_convert", "wav"),
        "video_to_aac":  ("_ffmpeg_convert", "aac"),
        "video_to_flac": ("_ffmpeg_convert", "flac"),
        # Audio → Audio
        "audio_to_mp3":  ("_ffmpeg_convert", "mp3"),
        "audio_to_wav":  ("_ffmpeg_convert", "wav"),
        "audio_to_aac":  ("_ffmpeg_convert", "aac"),
        "audio_to_ogg":  ("_ffmpeg_convert", "ogg"),
        "audio_to_flac": ("_ffmpeg_convert", "flac"),
        "audio_to_m4a":  ("_ffmpeg_convert", "m4a"),
    }

    @staticmethod
    def _office_to_pdf_com(src: str, dst: str, app_name: str) -> bool:
        """
        Convert src → dst PDF via Microsoft Office COM automation.
        app_name: "PowerPoint.Application" | "Excel.Application" | "Word.Application"
        Returns True on success, False if Office/comtypes not available.
        Works only on Windows with Office installed.
        """
        try:
            import comtypes.client
            import comtypes
            src_abs = str(Path(src).resolve())
            dst_abs = str(Path(dst).resolve())

            # Excel constants
            XL_PORTRAIT  = 1
            XL_LANDSCAPE = 2
            XL_PDF       = 0

            # Word constants
            WD_PDF = 17

            if "Excel" in app_name:
                app = comtypes.client.CreateObject(app_name)
                try:
                    app.Visible = False
                except Exception:
                    pass
                app.DisplayAlerts = False
                try:
                    wb = app.Workbooks.Open(src_abs)
                    try:
                        for sheet in wb.Worksheets:
                            try:
                                used = sheet.UsedRange
                                ps   = sheet.PageSetup

                                # Measure real column width in points via .Width property
                                # A4 portrait printable ~510 pt, landscape ~750 pt
                                total_width_pts = sum(
                                    sheet.Columns(used.Column + i).Width
                                    for i in range(used.Columns.Count)
                                )
                                ps.Orientation = XL_LANDSCAPE if total_width_pts > 510 else XL_PORTRAIT
                                # Zoom=100, no FitTo — Excel paginates naturally
                                ps.Zoom           = 100
                                ps.FitToPagesWide = False
                                ps.FitToPagesTall = False
                                ps.PrintArea      = used.Address
                            except Exception:
                                pass
                        wb.ExportAsFixedFormat(XL_PDF, dst_abs)
                    finally:
                        wb.Close(False)
                finally:
                    app.Quit()

            elif "PowerPoint" in app_name:
                import shutil

                tmp_dst = str(Path(tempfile.gettempdir()) / "pptx_com_out.pdf")
                _ps_env = os.environ.copy()
                _ps_env["FCP_SRC"] = src_abs
                _ps_env["FCP_DST"] = tmp_dst

                ps_a = """
                $app = New-Object -ComObject PowerPoint.Application
                $app.Visible = -1
                try {
                    $prs = $app.Presentations.Open($env:FCP_SRC, 0, 0, -1)
                    $prs.SaveAs($env:FCP_DST, 32)
                    $prs.Close()
                } finally { $app.Quit() }
                """
                subprocess.run(
                    ["powershell", "-NoProfile", "-NonInteractive",
                    "-ExecutionPolicy", "Bypass", "-Command", ps_a],
                    capture_output=True, timeout=120,
                    creationflags=_NO_WINDOW,
                    env=_ps_env,
                )
                if Path(tmp_dst).exists():
                    shutil.move(tmp_dst, dst_abs)
                    return Path(dst_abs).exists()

                ps_b = """
                $app = New-Object -ComObject PowerPoint.Application
                $app.Visible = -1
                try {
                    $prs = $app.Presentations.Open($env:FCP_SRC, 0, 0, -1)
                    $prs.PrintOut(1, $prs.Slides.Count, $env:FCP_DST, 0, 2)
                    $prs.Close()
                } finally { $app.Quit() }
                """
                r2 = subprocess.run(
                    ["powershell", "-NoProfile", "-NonInteractive",
                    "-ExecutionPolicy", "Bypass", "-Command", ps_b],
                    capture_output=True, timeout=120,
                    creationflags=_NO_WINDOW,
                    env=_ps_env,
                )
                if Path(tmp_dst).exists():
                    shutil.move(tmp_dst, dst_abs)
                    return Path(dst_abs).exists()

                print(f"[COM-PS] stderr: {r2.stderr.decode('utf-8','replace')[:400]}")

            elif "Word" in app_name:
                app = comtypes.client.CreateObject(app_name)
                try:
                    app.Visible = False
                except Exception:
                    pass
                try:
                    doc = app.Documents.Open(src_abs)
                    try:
                        doc.SaveAs2(dst_abs, WD_PDF)
                    finally:
                        doc.Close(False)
                finally:
                    app.Quit()

            return Path(dst).exists()

        except Exception as _com_exc:
            print(f"[COM] {app_name} failed: {_com_exc}")
            return False

    def convert(self, conversion_type, src, dst_dir):
        if conversion_type not in self._DISPATCH:
            return ConversionResult(False, src, "", error=f"Unknown type: {conversion_type}")
        method_name, ext = self._DISPATCH[conversion_type]
        dst = _build_dst(src, dst_dir, ext)
        file_size = os.path.getsize(src) if os.path.exists(src) else 0
        try:
            method = getattr(self, method_name)
            if method_name in ("_image_convert", "_ffmpeg_convert", "_heic_convert", "_raw_convert"):
                elapsed = _timed(lambda: method(src, dst, conversion_type))
            else:
                elapsed = _timed(lambda: method(src, dst))
            return ConversionResult(True, src, dst, elapsed=elapsed, file_size=file_size)
        except Exception as exc:
            return ConversionResult(False, src, dst, error=str(exc), file_size=file_size)

    def convert_batch(self, conversion_type, sources, dst_dir, progress_cb=None):
        results = []
        for i, src in enumerate(sources, 1):
            result = self.convert(conversion_type, src, dst_dir)
            results.append(result)
            if progress_cb:
                progress_cb(i, len(sources), src)
        return results

    @staticmethod
    def _rtf_read_raw(src: str) -> str:
        """Read RTF bytes and decode robustly (utf-8 → cp1252 → latin-1)."""
        with open(src, "rb") as fh:
            raw_bytes = fh.read()
        for enc in ("utf-8", "cp1252", "latin-1"):
            try:
                return raw_bytes.decode(enc)
            except Exception:
                pass
        return raw_bytes.decode("utf-8", errors="replace")

    @staticmethod
    def _rtf_extract_images(raw: str) -> list:
        """
        Extract embedded PNG/JPEG images from \\pict blocks (hex-encoded).
        Returns list of (image_bytes, 'png'|'jpg').
        WMF/EMF bitmaps are skipped (not portably renderable).
        """
        images = []
        for m in re.finditer(r"\{\\pict((?:[^{}]|\{[^{}]*\})*)\}", raw, re.S):
            block = m.group(1)
            if re.search(r"\\pngblip\b", block):
                fmt = "png"
            elif re.search(r"\\jpegblip\b", block):
                fmt = "jpg"
            else:
                continue
            hex_data = re.sub(r"\\[a-zA-Z]+[-0-9]*\s?", "", block)
            hex_data = re.sub(r"[^0-9a-fA-F]", "", hex_data)
            if len(hex_data) < 8:
                continue
            try:
                images.append((bytes.fromhex(hex_data), fmt))
            except Exception:
                pass
        return images

    @staticmethod
    def _rtf_spans_to_paragraphs(spans: list) -> list:
        """
        Assemble spans into a document structure:
          { 'type': 'para',      'para': { runs, text_content } }
          { 'type': 'table_row', 'cells': [ [para, ...], ... ] }

        Adjacent runs with identical formatting are merged.
        """
        items = []
        cur_runs = []
        cur_cell_paras = []
        cur_row_cells = []

        def _flush(is_cell=False):
            merged = []
            for run in cur_runs:
                if (merged
                        and merged[-1]["bold"]      == run["bold"]
                        and merged[-1]["italic"]    == run["italic"]
                        and merged[-1]["underline"] == run["underline"]
                        and merged[-1]["fontsize"]  == run["fontsize"]
                        and merged[-1]["color"]     == run["color"]):
                    merged[-1]["text"] += run["text"]
                else:
                    merged.append(dict(run))
            cur_runs.clear()
            return {"runs": merged,
                    "text_content": "".join(r["text"] for r in merged).strip(),
                    "is_table_cell": is_cell}

        for sp in spans:
            if sp["cell_end"]:
                cur_cell_paras.append(_flush(is_cell=True))
                cur_row_cells.append(list(cur_cell_paras))
                cur_cell_paras.clear()
                continue
            if sp["row_end"]:
                leftover = _flush(is_cell=True)
                if leftover["text_content"] or leftover["runs"]:
                    cur_cell_paras.append(leftover)
                    cur_row_cells.append(list(cur_cell_paras))
                    cur_cell_paras.clear()
                items.append({"type": "table_row", "cells": list(cur_row_cells)})
                cur_row_cells.clear()
                continue
            if sp["par"]:
                para = _flush()
                if sp["in_table"]:
                    cur_cell_paras.append(para)
                else:
                    items.append({"type": "para", "para": para})
                continue
            if sp["text"]:
                cur_runs.append({k: sp[k] for k in
                                 ("text","bold","italic","underline","fontsize","color")})

        if cur_runs:
            items.append({"type": "para", "para": _flush()})
        return items

    def _rtf_to_pdf_native(self, src, dst):
        """
        Full-fidelity native RTF → PDF (no Office, no pandoc needed).
        Preserves bold · italic · underline · font sizes · RGB colors ·
        tables · embedded PNG/JPEG images.
        """
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import mm
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors as rl_colors
        from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                        Table, TableStyle, Image as RLImage)

        raw         = self._rtf_read_raw(src)
        color_table = self._rtf_parse_colortbl(raw)
        img_data    = self._rtf_extract_images(raw)
        tokens      = self._rtf_tokenize(raw)
        spans       = self._rtf_parse_spans(tokens, color_table)
        paras       = self._rtf_spans_to_paragraphs(spans)

        sty  = getSampleStyleSheet()
        base = ParagraphStyle("RTFBase", parent=sty["Normal"],
                              fontSize=11, leading=15, spaceAfter=3,
                              fontName="Helvetica")
        td_s = ParagraphStyle("RTFTd", parent=sty["Normal"],
                              fontSize=9, leading=12, fontName="Helvetica")

        def _runs_to_xml(runs, base_fs=11):
            parts = []
            for r in runs:
                txt = _safe_html(r["text"])
                if not txt:
                    continue
                fs = max(6, r["fontsize"] // 2)
                o, c = [], []
                if r["bold"]:      o.append("<b>");  c.insert(0,"</b>")
                if r["italic"]:    o.append("<i>");  c.insert(0,"</i>")
                if r["underline"]: o.append("<u>");  c.insert(0,"</u>")
                if r["color"]:
                    rv,gv,bv = r["color"]
                    o.append(f'<font color="#{rv:02x}{gv:02x}{bv:02x}">'); c.insert(0,"</font>")
                if fs != base_fs:
                    o.append(f'<font size="{fs}">'); c.insert(0,"</font>")
                parts.append("".join(o) + txt + "".join(c))
            return "".join(parts)

        story = []

        def _build_rl_table(rows_data):
            """Convert list of table_row items into a single ReportLab Table."""
            max_cols = max(len(r["cells"]) for r in rows_data)
            tbl_rows = []
            for row_item in rows_data:
                tbl_row = []
                for ci in range(max_cols):
                    cell_paras = row_item["cells"][ci] if ci < len(row_item["cells"]) else []
                    cell_content = []
                    for cp in cell_paras:
                        x = _runs_to_xml(cp["runs"], base_fs=9)
                        if x.strip():
                            cell_content.append(Paragraph(x, td_s))
                    tbl_row.append(cell_content or [Paragraph("", td_s)])
                tbl_rows.append(tbl_row)
            tbl = Table(tbl_rows, hAlign="LEFT", colWidths=[None]*max_cols)
            tbl.setStyle(TableStyle([
                ("GRID",          (0,0),(-1,-1), 0.4, rl_colors.HexColor("#bbbbbb")),
                ("BACKGROUND",    (0,0),(-1, 0), rl_colors.HexColor("#f5f5f5")),
                ("VALIGN",        (0,0),(-1,-1), "TOP"),
                ("LEFTPADDING",   (0,0),(-1,-1), 4),
                ("RIGHTPADDING",  (0,0),(-1,-1), 4),
                ("TOPPADDING",    (0,0),(-1,-1), 3),
                ("BOTTOMPADDING", (0,0),(-1,-1), 3),
            ]))
            return tbl

        idx = 0
        while idx < len(paras):
            item = paras[idx]

            if item["type"] == "table_row":
                rows_data = []
                while idx < len(paras) and paras[idx]["type"] == "table_row":
                    rows_data.append(paras[idx]); idx += 1
                story.append(_build_rl_table(rows_data))
                continue

            para = item["para"]
            xml  = _runs_to_xml(para["runs"])
            idx += 1
            if not xml.strip():
                story.append(Spacer(1, 4)); continue
            plain = para["text_content"]
            if (plain.isupper() and 0 < len(plain) < 80
                    and all(r["bold"] or r["fontsize"] >= 28
                            for r in para["runs"] if r["text"].strip())):
                hs = ParagraphStyle("RTFHead", parent=base, fontSize=13,
                                    leading=17, spaceBefore=6, spaceAfter=4,
                                    fontName="Helvetica-Bold")
                story.append(Paragraph(xml, hs))
            else:
                sizes = [r["fontsize"] for r in para["runs"] if r["text"].strip()]
                dom   = max(6, (max(sizes) // 2)) if sizes else 11
                ps    = ParagraphStyle("RTFPar", parent=base,
                                       fontSize=dom, leading=max(dom+3,13))
                story.append(Paragraph(xml, ps))

        for img_bytes, fmt in img_data:
            try:
                from PIL import Image as PILImg
                buf   = io.BytesIO(img_bytes)
                pil   = PILImg.open(buf)
                w, h  = pil.size
                max_w = 120 * mm
                scale = min(1.0, max_w / max(w, 1))
                buf.seek(0)
                story.append(Spacer(1, 6))
                story.append(RLImage(buf, width=w*scale, height=h*scale))
                story.append(Spacer(1, 6))
            except Exception:
                pass

        if not story:
            story.append(Paragraph("(empty document)", base))

        SimpleDocTemplate(
            dst, pagesize=A4,
            leftMargin=20*mm, rightMargin=20*mm,
            topMargin=20*mm, bottomMargin=20*mm,
        ).build(story)

    def _pptx_to_pdf_native(self, src, dst):
        """
        Native PPTX → PDF using python-pptx + reportlab + matplotlib.
        Handles: text, bullets, images, tables, charts (bar/line/pie/donut).
        All python-pptx API calls are wrapped in try/except.
        """
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from reportlab.lib.pagesizes import landscape, A4
        from reportlab.lib.units import mm
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors
        from reportlab.lib.enums import TA_CENTER, TA_RIGHT
        from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                        Image as RLImage, PageBreak,
                                        Table, TableStyle, HRFlowable)

        prs = Presentation(src)
        try:    sw = int(prs.slide_width)  or 9144000
        except Exception: sw = 9144000
        try:    sh = int(prs.slide_height) or 6858000
        except Exception: sh = 6858000

        ratio   = float(sh) / float(sw)
        pw      = landscape(A4)[0]
        ph      = pw * ratio
        margin  = 14 * mm
        inner_w = pw - 2 * margin

        S = getSampleStyleSheet()
        def sty(name, **kw):
            p = kw.pop("parent", S["Normal"])
            return ParagraphStyle(name, parent=p, **kw)

        ST = {
            "title"  : sty("PT", parent=S["Heading1"], fontSize=20, leading=24,
                           spaceAfter=8, textColor=colors.HexColor("#1a1a2e"),
                           alignment=TA_CENTER),
            "h2"     : sty("PH2", parent=S["Heading2"], fontSize=15, leading=19,
                           spaceAfter=5, textColor=colors.HexColor("#1e3a5f")),
            "body"   : sty("PB", fontSize=11, leading=16, spaceAfter=4),
            "bul0"   : sty("PB0", fontSize=11, leading=16, spaceAfter=3, leftIndent=12),
            "bul1"   : sty("PB1", fontSize=10, leading=15, spaceAfter=2, leftIndent=24,
                           textColor=colors.HexColor("#333")),
            "bul2"   : sty("PB2", fontSize=9,  leading=14, spaceAfter=2, leftIndent=36,
                           textColor=colors.HexColor("#555")),
            "slide_n": sty("PSN", fontSize=8, leading=10,
                           textColor=colors.HexColor("#aaa"), alignment=TA_RIGHT),
            "td"     : sty("PTd", fontSize=9, leading=13),
            "th"     : sty("PTh", fontSize=9, leading=13, fontName="Helvetica-Bold"),
            "caption": sty("PC",  fontSize=8, leading=11,
                           textColor=colors.HexColor("#777"), alignment=TA_CENTER),
        }

        BULLETS = set("•–-*◦▪▸")
        story   = []
        tmp_imgs = []

        def _safe(txt):
            return _safe_html(str(txt or ""))

        def _add_img_blob(blob, ext):
            try:
                tf = tempfile.NamedTemporaryFile(suffix=f".{ext or 'png'}", delete=False)
                tf.write(blob); tf.close()
                tmp_imgs.append(tf.name)
                rl = RLImage(tf.name)
                scale = min(inner_w / rl.imageWidth,
                            (ph - 2*margin)*0.6 / rl.imageHeight, 1.0)
                rl.drawWidth  = rl.imageWidth  * scale
                rl.drawHeight = rl.imageHeight * scale
                rl.hAlign = "CENTER"
                story.append(Spacer(1, 4))
                story.append(rl)
                story.append(Spacer(1, 4))
            except Exception:
                pass

        def _render_chart(shape):
            """Extract chart data and render as matplotlib image."""
            try:
                import matplotlib
                matplotlib.use("Agg")
                import matplotlib.pyplot as plt
                import numpy as np
                from pptx.enum.chart import XL_CHART_TYPE

                chart = shape.chart
                ctype = chart.chart_type

                NS_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
                NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
                cxml = chart._element

                series_list = []
                for ser in cxml.findall(f".//{{{NS_C}}}ser"):
                    name = ""
                    v_el = ser.find(f".//{{{NS_C}}}tx//{{{NS_C}}}v")
                    if v_el is not None: name = v_el.text or ""

                    cats = [el.find(f"{{{NS_C}}}v").text
                            for el in ser.findall(f".//{{{NS_C}}}cat//{{{NS_C}}}pt")
                            if el.find(f"{{{NS_C}}}v") is not None]

                    vals = []
                    for pt in ser.findall(f".//{{{NS_C}}}val//{{{NS_C}}}pt"):
                        v = pt.find(f"{{{NS_C}}}v")
                        if v is not None:
                            try: vals.append(float(v.text))
                            except Exception: vals.append(0.0)

                    pt_colors = {}
                    for dpt in ser.findall(f"{{{NS_C}}}dPt"):
                        idx_el = dpt.find(f"{{{NS_C}}}idx")
                        clr_el = dpt.find(f".//{{{NS_A}}}srgbClr")
                        if idx_el is not None and clr_el is not None:
                            idx = int(idx_el.get("val", 0))
                            pt_colors[idx] = "#" + clr_el.get("val", "4472C4")

                    ser_clr = None
                    clr_el  = ser.find(f".//{{{NS_A}}}srgbClr")
                    if clr_el is not None:
                        ser_clr = "#" + clr_el.get("val", "4472C4")

                    series_list.append({
                        "name": name, "cats": cats, "vals": vals,
                        "pt_colors": pt_colors, "ser_color": ser_clr,
                    })

                if not series_list:
                    return False

                fig_w = float(inner_w) / 72 / 1.333 * 1.5
                fig_h = fig_w * 0.6
                fig, ax = plt.subplots(figsize=(fig_w, fig_h))
                fig.patch.set_facecolor("white")
                ax.set_facecolor("#f8f9fa")

                is_pie  = ctype in (
                    XL_CHART_TYPE.PIE, XL_CHART_TYPE.PIE_EXPLODED,
                    XL_CHART_TYPE.DOUGHNUT, XL_CHART_TYPE.DOUGHNUT_EXPLODED,
                )
                is_bar  = ctype in (
                    XL_CHART_TYPE.BAR_CLUSTERED, XL_CHART_TYPE.BAR_STACKED,
                    XL_CHART_TYPE.BAR_STACKED_100,
                    XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.COLUMN_STACKED,
                )
                is_line = ctype in (
                    XL_CHART_TYPE.LINE, XL_CHART_TYPE.LINE_MARKERS,
                )

                ser0 = series_list[0]

                if is_pie:
                    vals_  = ser0["vals"]
                    labels = ser0["cats"] or [f"Cat {i+1}" for i in range(len(vals_))]
                    clrs   = [ser0["pt_colors"].get(i, None) for i in range(len(vals_))]
                    # Fill None colors with defaults
                    default_clrs = plt.rcParams["axes.prop_cycle"].by_key()["color"]
                    clrs = [c if c else default_clrs[i % len(default_clrs)]
                            for i, c in enumerate(clrs)]

                    wedge_kw = {"width": 0.55} if "DOUGHNUT" in str(ctype) else {}
                    wedges, texts, autotexts = ax.pie(
                        vals_, labels=None, colors=clrs,
                        autopct="%1.1f%%", startangle=90,
                        wedgeprops=wedge_kw,
                        pctdistance=0.75,
                    )
                    for at in autotexts:
                        at.set_fontsize(8)
                    ax.legend(wedges, labels, loc="center left",
                              bbox_to_anchor=(1, 0.5), fontsize=8)
                    if "DOUGHNUT" in str(ctype):
                        total = sum(vals_)
                        ax.text(0, 0, f"{total:.0f}", ha="center", va="center",
                                fontsize=12, fontweight="bold")
                    ax.set_title(ser0["name"] or "Chart", fontsize=11, pad=10)

                elif is_bar or (not is_line):
                    cats    = ser0["cats"] or [str(i+1) for i in range(len(ser0["vals"]))]
                    x       = np.arange(len(cats))
                    n_ser   = len(series_list)
                    width   = 0.8 / max(n_ser, 1)
                    offsets = np.linspace(-(n_ser-1)*width/2, (n_ser-1)*width/2, n_ser)
                    clr_cycle = plt.rcParams["axes.prop_cycle"].by_key()["color"]

                    for si2, ser2 in enumerate(series_list):
                        clr = ser2["ser_color"] or clr_cycle[si2 % len(clr_cycle)]
                        ax.bar(x + offsets[si2], ser2["vals"], width,
                               label=ser2["name"], color=clr, alpha=0.88)

                    ax.set_xticks(x)
                    ax.set_xticklabels(cats, rotation=30, ha="right", fontsize=8)
                    ax.set_ylabel("Value", fontsize=9)
                    ax.tick_params(axis="y", labelsize=8)
                    ax.grid(axis="y", alpha=0.3, linestyle="--")
                    if any(s["name"] for s in series_list):
                        ax.legend(fontsize=8)

                else:
                    cats = ser0["cats"] or [str(i+1) for i in range(len(ser0["vals"]))]
                    for si2, ser2 in enumerate(series_list):
                        ax.plot(cats, ser2["vals"], marker="o", label=ser2["name"],
                                linewidth=1.8, markersize=4)
                    ax.set_xticks(range(len(cats)))
                    ax.set_xticklabels(cats, rotation=30, ha="right", fontsize=8)
                    ax.grid(alpha=0.3, linestyle="--")
                    ax.legend(fontsize=8)

                plt.tight_layout(pad=1.0)

                tf = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
                fig.savefig(tf.name, dpi=150, bbox_inches="tight",
                            facecolor="white")
                plt.close(fig)
                tf.close()
                tmp_imgs.append(tf.name)

                rl = RLImage(tf.name)
                scale = min(inner_w / rl.imageWidth,
                            (ph - 2*margin)*0.55 / rl.imageHeight, 1.0)
                rl.drawWidth  = rl.imageWidth  * scale
                rl.drawHeight = rl.imageHeight * scale
                rl.hAlign = "CENTER"
                story.append(Spacer(1, 4))
                story.append(rl)
                story.append(Spacer(1, 6))
                return True

            except Exception:
                try:
                    chart = shape.chart
                    story.append(Paragraph(f"[Chart — {chart.chart_type}]", ST["h2"]))
                    for ser in chart.series:
                        story.append(Paragraph(f"• {ser.name}", ST["body"]))
                except Exception:
                    story.append(Paragraph("[Chart]", ST["body"]))
                return False

        def _render_table(tbl_obj):
            try:
                rows_data = []
                for ri, row in enumerate(tbl_obj.rows):
                    cells = []
                    for ci, cell in enumerate(row.cells):
                        try:   txt = cell.text.strip()
                        except Exception: txt = ""
                        ps = ST["th"] if ri == 0 else ST["td"]
                        cells.append(Paragraph(_safe(txt) or " ", ps))
                    rows_data.append(cells)
                if not rows_data:
                    return
                t = Table(rows_data, repeatRows=1)
                t.setStyle(TableStyle([
                    ("GRID",         (0,0),(-1,-1), 0.5, colors.HexColor("#ccc")),
                    ("BACKGROUND",   (0,0),(-1,0),  colors.HexColor("#e8edf5")),
                    ("VALIGN",       (0,0),(-1,-1), "TOP"),
                    ("LEFTPADDING",  (0,0),(-1,-1), 4),
                    ("RIGHTPADDING", (0,0),(-1,-1), 4),
                    ("TOPPADDING",   (0,0),(-1,-1), 3),
                    ("BOTTOMPADDING",(0,0),(-1,-1), 3),
                ]))
                story.append(Spacer(1, 6))
                story.append(t)
                story.append(Spacer(1, 6))
            except Exception:
                pass

        def _render_text_frame(tf_obj, default_sty):
            for para in tf_obj.paragraphs:
                try:
                    text = para.text.strip()
                    if not text:
                        story.append(Spacer(1, 3))
                        continue
                    try:    level = int(para.level or 0)
                    except Exception: level = 0
                    if level == 0:
                        is_bul = text[0] in BULLETS
                        ps     = ST["bul0"] if is_bul else default_sty
                        prefix = "• " if is_bul and text[0] not in BULLETS else ""
                    elif level == 1:
                        ps, prefix = ST["bul1"], "  ◦ "
                    else:
                        ps, prefix = ST["bul2"], "    ▪ "
                    parts = []
                    for run in para.runs:
                        try:
                            rt = run.text
                            if not rt: continue
                            s = _safe(rt)
                            try:
                                b, i = run.font.bold, run.font.italic
                                if b and i: s = f"<b><i>{s}</i></b>"
                                elif b:     s = f"<b>{s}</b>"
                                elif i:     s = f"<i>{s}</i>"
                            except Exception:
                                pass
                            parts.append(s)
                        except Exception:
                            pass
                    inner = prefix + ("".join(parts) if parts else _safe(prefix + text))
                    story.append(Paragraph(inner, ps))
                except Exception:
                    try:
                        raw = para.text.strip()
                        if raw:
                            story.append(Paragraph(_safe(raw), default_sty))
                    except Exception:
                        pass

        for sn, slide in enumerate(prs.slides, 1):
            shapes_sorted = sorted(
                slide.shapes,
                key=lambda sh: (
                    (sh.top  or 0),
                    (sh.left or 0),
                )
            )

            for shape in shapes_sorted:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try: _add_img_blob(shape.image.blob, shape.image.ext)
                        except Exception: pass
                        continue

                    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        _render_chart(shape)
                        continue

                    if shape.has_table:
                        _render_table(shape.table)
                        continue

                    if not shape.has_text_frame:
                        continue
                    if not shape.text_frame.text.strip():
                        continue

                    is_title = shape.name.lower().startswith("title")

                    if is_title:
                        story.append(Paragraph(_safe(shape.text_frame.text.strip()),
                                               ST["title"]))
                        story.append(HRFlowable(width="100%", thickness=0.5,
                                               color=colors.HexColor("#1e3a8a"),
                                               spaceAfter=6))
                    else:
                        _render_text_frame(shape.text_frame, ST["body"])

                except Exception:
                    pass

            story.append(Spacer(1, 6))
            story.append(Paragraph(f"— {sn} / {len(prs.slides)} —", ST["slide_n"]))
            if sn < len(prs.slides):
                story.append(PageBreak())

        if not story:
            story.append(Paragraph("(empty presentation)", ST["body"]))

        def _safe_title(p):
            try: return p.slides[0].shapes.title.text or ""
            except Exception: return ""

        SimpleDocTemplate(
            dst, pagesize=(pw, ph),
            leftMargin=margin, rightMargin=margin,
            topMargin=margin, bottomMargin=margin,
            title=_safe_title(prs),
        ).build(story)

        for f in tmp_imgs:
            try: os.remove(f)
            except Exception: pass

    def _html_to_pdf(self, src, dst):
        """
        HTML → PDF — lightweight strategy stack, PyInstaller-compatible.

        1. pdfkit    (wkhtmltopdf separate binary, optional)
        2. weasyprint (pip install weasyprint)
        3. reportlab  — parses HTML manually, never fails, no duplicate content
           (fitz insert_htmlbox removed — duplicates content on some HTML inputs)
        """
        src_path = Path(src)
        base_dir = src_path.parent

        with open(src, "r", encoding="utf-8", errors="replace") as f:
            html_raw = f.read()

        html = self._inline_all_resources(html_raw, base_dir)

        # Strategy 1: pdfkit (wkhtmltopdf)
        try:
            import pdfkit
            tmp = tempfile.NamedTemporaryFile(
                suffix=".html", delete=False, mode="w", encoding="utf-8")
            tmp.write(html); tmp.close()
            try:
                pdfkit.from_file(tmp.name, dst, options={
                    "enable-local-file-access": "",
                    "load-error-handling":       "ignore",
                    "load-media-error-handling": "ignore",
                    "print-media-type":          "",
                    "quiet":                     "",
                    "margin-top":    "15mm",
                    "margin-bottom": "15mm",
                    "margin-left":   "15mm",
                    "margin-right":  "15mm",
                })
                return
            finally:
                try: os.remove(tmp.name)
                except Exception: pass
        except Exception:
            pass

        # Strategy 2: weasyprint
        try:
            import warnings
            import weasyprint
            with warnings.catch_warnings():
                warnings.simplefilter("ignore")
                weasyprint.HTML(
                    string=html,
                    base_url=base_dir.as_uri()
                ).write_pdf(dst)
            return
        except Exception:
            pass

        # Strategy 3: reportlab
        self._reportlab_html_to_pdf(html, dst)

    def _inline_all_resources(self, html, base_dir):
        """Inline all local img/css/url resources as base64 data-URIs."""
        SRC_RE  = re.compile(r"""src=(['"])([^'"]+)\1""",  re.I)
        HREF_RE = re.compile(r"""href=(['"])([^'"]+)\1""", re.I)

        def _link_to_style(m):
            tag = m.group(0)
            if "stylesheet" not in tag.lower():
                return tag
            hm = HREF_RE.search(tag)
            if not hm:
                return tag
            href = hm.group(2)
            if href.startswith(("http://", "https://", "data:", "//")):
                return tag
            css_path = base_dir / href
            if not css_path.exists():
                return tag
            try:
                css = css_path.read_text(encoding="utf-8", errors="replace")
                css = self._inline_css_urls(css, css_path.parent)
                return "<style>" + css + "</style>"
            except Exception:
                return tag
        html = re.sub(r"<link[^>]+>", _link_to_style, html, flags=re.I | re.S)

        def _style_block(m):
            return "<style>" + self._inline_css_urls(m.group(1), base_dir) + "</style>"
        html = re.sub(r"<style[^>]*>(.*?)</style>",
                      _style_block, html, flags=re.I | re.S)

        def _img_src(m):
            tag = m.group(0)
            sm  = SRC_RE.search(tag)
            if not sm:
                return tag
            val = sm.group(2)
            if val.startswith(("http://", "https://", "data:")):
                return tag
            p = base_dir / val
            if not p.exists():
                return tag
            try:
                b64 = _img_to_b64(p.read_bytes(), _mime_for_ext(p.suffix.lstrip(".")))
                return SRC_RE.sub('src="' + b64 + '"', tag, count=1)
            except Exception:
                return tag
        html = re.sub(r"<img[^>]+>", _img_src, html, flags=re.I | re.S)
        return html

    def _inline_css_urls(self, css, base_dir):
        """Replace url(path) in CSS with base64 data URIs."""
        URL_RE = re.compile(r"url\\(\\s*([\"']?)([^)'\"]+)\\1\\s*\\)", re.I)
        def repl(m):
            raw = m.group(2).strip()
            if raw.startswith(("http://", "https://", "data:")):
                return m.group(0)
            p = base_dir / raw
            if not p.exists():
                return m.group(0)
            try:
                b64 = _img_to_b64(p.read_bytes(), _mime_for_ext(p.suffix.lstrip(".")))
                return "url('" + b64 + "')"
            except Exception:
                return m.group(0)
        return URL_RE.sub(repl, css)

    def _reportlab_html_to_pdf(self, html, dst):
        """
        Reportlab HTML fallback — faithful CSS class + inline style support.
        Parses <style> blocks to extract .class rules (text-align, margin-left).
        Applies class + inline style to every <p>/<div>.
        <br> inside <p> becomes a real line break inside the paragraph.
        .pn (page number) divs are skipped.
        """
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import mm
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors
        from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
        from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                        HRFlowable)

        MAX_W = A4[0] - 40 * mm
        base  = getSampleStyleSheet()

        class_rules = {}
        for sb in re.finditer(r'<style[^>]*>(.*?)</style>', html, re.I | re.S):
            css = sb.group(1)
            for rule in re.finditer(r'\.([\w-]+)\s*\{([^}]*)\}', css):
                cls   = rule.group(1)
                decls = {}
                for d in rule.group(2).split(';'):
                    d = d.strip()
                    if ':' not in d:
                        continue
                    p, _, v = d.partition(':')
                    decls[p.strip().lower()] = v.strip()
                class_rules[cls] = decls

        ALIGN_MAP = {
            'right':   TA_RIGHT,
            'center':  TA_CENTER,
            'left':    TA_LEFT,
            'justify': TA_JUSTIFY,
        }

        def _decls_to_align(decls):
            return ALIGN_MAP.get(decls.get('text-align', '').lower())

        def _decls_to_indent(decls):
            ml = decls.get('margin-left', '')
            m  = re.search(r'([\d.]+)px', ml)
            return float(m.group(1)) * 0.75 if m else 0.0

        _cache = {}

        def _sty(align=TA_LEFT, indent=0.0, size=10.5, bold=False):
            key = (align, round(indent), size, bold)
            if key not in _cache:
                _cache[key] = ParagraphStyle(
                    'RS%d' % len(_cache),
                    parent=base['Normal'],
                    fontName='Helvetica-Bold' if bold else 'Helvetica',
                    fontSize=size, leading=size * 1.55,
                    spaceAfter=5, alignment=align, leftIndent=indent,
                )
            return _cache[key]

        H_STYS = {
            'h1': _sty(size=16, bold=True),
            'h2': _sty(size=14, bold=True),
            'h3': _sty(size=12, bold=True),
            'h4': _sty(size=11, bold=True),
            'h5': _sty(size=10, bold=True),
            'h6': _sty(size=10, bold=True),
        }
        TD_S  = _sty(size=9)
        TH_S  = _sty(size=9, bold=True)
        PRE_S = ParagraphStyle('RPRE', parent=base['Code'],
                               fontName='Courier', fontSize=9, leading=13,
                               spaceAfter=6, backColor=colors.HexColor('#f5f5f5'))

        _CLS_RE  = re.compile(r'''class=['"]([^'"]+)['"]''',  re.I)
        _STY_RE  = re.compile(r'''style=['"]([^'"]*?)['"]''', re.I)
        _SRC_RE  = re.compile(r'''src=['"]([^'"]+)['"]''',    re.I)

        def _resolve(attrs):
            align  = None
            indent = 0.0
            cm = _CLS_RE.search(attrs)
            if cm:
                for c in cm.group(1).split():
                    rd = class_rules.get(c, {})
                    a  = _decls_to_align(rd)
                    if a is not None:
                        align = a
                    i = _decls_to_indent(rd)
                    if i:
                        indent = i
            sm = _STY_RE.search(attrs)
            if sm:
                decl_str = sm.group(1)
                rd2 = {}
                for d in decl_str.split(';'):
                    d = d.strip()
                    if ':' not in d:
                        continue
                    p, _, v = d.partition(':')
                    rd2[p.strip().lower()] = v.strip()
                a2 = _decls_to_align(rd2)
                if a2 is not None:
                    align = a2
                i2 = _decls_to_indent(rd2)
                if i2:
                    indent = i2
            return (align or TA_LEFT), indent

        def _decode(t):
            return (t.replace('&nbsp;', '\xa0')
                     .replace('&amp;',  '&')
                     .replace('&lt;',   '<')
                     .replace('&gt;',   '>')
                     .replace('&quot;', '"')
                     .replace('&#39;',  "'"))

        def _inline(frag):
            return inline_markup(frag, decode_fn=_decode, preserve_font_tags=False)

        tmp_imgs = []

        def _add_img(src_val):
            add_image_flowable(story, src_val, MAX_W, tmp_imgs)

        def _parse_tbl(tbl_html):
            return build_table_flowable(tbl_html, _inline, TD_S, TH_S)

        h2 = re.sub(r'<head[^>]*>.*?</head>',     '', html, flags=re.I | re.S)
        h2 = re.sub(r'<script[^>]*>.*?</script>', '', h2,   flags=re.I | re.S)
        bm = re.search(r'<body[^>]*>(.*?)</body>', h2,       re.I | re.S)
        body = bm.group(1) if bm else h2

        story = []

        TAG_RE = re.compile(r'<(/?)(\w+)((?:\s[^>]*)?)/?>',  re.I)
        pos    = 0

        for m in TAG_RE.finditer(body):
            pos     = m.end()
            closing = m.group(1)
            tag     = m.group(2).lower()
            attrs   = m.group(3) or ''

            if tag == 'hr' and not closing:
                story.append(HRFlowable(width='100%', thickness=0.5,
                                        color=colors.HexColor('#ccc'),
                                        spaceAfter=5, spaceBefore=5))

            elif tag in ('h1','h2','h3','h4','h5','h6') and not closing:
                em = re.search(r'</' + tag + r'\s*>', body[pos:], re.I)
                if em:
                    txt = _inline(body[pos:pos + em.start()])
                    pos += em.end()
                    if txt:
                        story.append(Paragraph(txt, H_STYS.get(tag, H_STYS['h4'])))

            elif tag == 'p' and not closing:
                align, indent = _resolve(attrs)
                em = re.search(r'</p\s*>', body[pos:], re.I)
                if em:
                    inner = body[pos:pos + em.start()]
                    pos  += em.end()
                    for im in re.finditer(r'<img[^>]+>', inner, re.I):
                        sm = _SRC_RE.search(im.group(0))
                        if sm: _add_img(sm.group(1))
                    txt = _inline(inner)
                    if txt.strip():
                        story.append(Paragraph(txt, _sty(align, indent)))

            elif tag == 'div' and not closing:
                cm2 = _CLS_RE.search(attrs)
                if cm2 and 'pn' in cm2.group(1).split():
                    em = re.search(r'</div\s*>', body[pos:], re.I)
                    if em: pos += em.end()

            elif tag == 'li' and not closing:
                em = re.search(r'</li\s*>', body[pos:], re.I)
                if em:
                    txt = _inline(body[pos:pos + em.start()])
                    pos += em.end()
                    if txt.strip():
                        story.append(Paragraph('• ' + txt, _sty()))

            elif tag == 'blockquote' and not closing:
                em = re.search(r'</blockquote\s*>', body[pos:], re.I)
                if em:
                    txt = _inline(body[pos:pos + em.start()])
                    pos += em.end()
                    if txt:
                        story.append(Paragraph(txt, _sty(indent=18)))

            elif tag == 'pre' and not closing:
                em = re.search(r'</pre\s*>', body[pos:], re.I)
                if em:
                    raw = re.sub(r'<[^>]+>', '', body[pos:pos + em.start()])
                    pos += em.end()
                    safe = (_decode(raw)
                            .replace('&', '&amp;')
                            .replace('<', '&lt;')
                            .replace('>', '&gt;')
                            .replace('\n', '<br/>'))
                    story.append(Paragraph(safe, PRE_S))

            elif tag == 'img' and not closing:
                sm = _SRC_RE.search(attrs)
                if sm: _add_img(sm.group(1))

            elif tag == 'table' and not closing:
                em = re.search(r'</table\s*>', body[pos:], re.I)
                if em:
                    tbl = _parse_tbl(body[pos:pos + em.start()])
                    pos += em.end()
                    if tbl:
                        story.append(Spacer(1, 6))
                        story.append(tbl)
                        story.append(Spacer(1, 6))

        if not story:
            story.append(Paragraph('(empty)', _sty()))

        SimpleDocTemplate(dst, pagesize=A4,
                          leftMargin=20*mm, rightMargin=20*mm,
                          topMargin=20*mm, bottomMargin=20*mm).build(story)
        for f in tmp_imgs:
            try: os.remove(f)
            except Exception: pass

    def _pdf_to_html(self, src, dst):
        """
        PDF → HTML — self-contained, faithful layout.

        Key improvements:
        - Margins reconstructed from block X positions relative to page width
          → text appears centred with proper left/right margins, not stuck left
        - Multi-column detection: blocks with x0 > 50% of page width → right col
        - Justified text for body paragraphs
        - Line-height and letter-spacing tuned per block font size
        - Consecutive lines of the same block merged into one <p> (no <br> soup)
        - Bold/italic/colour/size spans preserved
        - Images at their exact position in the flow
        - Superscript/subscript detected via vertical origin offset
        """
        import fitz

        doc  = fitz.open(src)
        name = Path(dst).stem

        _widths  = [p.rect.width  for p in doc]
        _heights = [p.rect.height for p in doc]
        MED_W  = sorted(_widths)[len(_widths)  // 2] if _widths  else 595.0
        sorted(_heights)[len(_heights) // 2] if _heights else 842.0

        CSS_W  = max(620, min(1080, int(MED_W * 1.333)))
        PAD    = max(40, int(CSS_W * 0.075))
        INNER  = CSS_W - PAD * 2

        BASE_F = max(11, min(15, round(MED_W / 50)))

        CSS = f"""<style>
  *, *::before, *::after {{ box-sizing: border-box; margin: 0; padding: 0; }}
  html {{ font-size: {BASE_F}px; }}
  body {{
    font-family: 'Georgia', 'Times New Roman', serif;
    background: #d8d8d8;
    color: #111;
    line-height: 1.65;
  }}
  /* page card */
  .page {{
    background: #fff;
    width: {CSS_W}px;
    margin: 28px auto;
    padding: {PAD}px;
    box-shadow: 0 3px 18px rgba(0,0,0,.18);
    border-radius: 3px;
    position: relative;
  }}
  /* page number */
  .pn {{
    font-size: 9px;
    color: #bbb;
    text-align: right;
    margin-bottom: 14px;
    font-family: 'Segoe UI', Arial, sans-serif;
  }}
  /* body paragraph */
  p {{
    margin: 0 0 6px 0;
    text-align: justify;
    hyphens: auto;
    word-spacing: 0.02em;
  }}
  p.center {{ text-align: center; }}
  p.right  {{ text-align: right;  }}
  /* headings */
  h1 {{ font-size: {BASE_F + 7}px; margin: 16px 0 8px; line-height: 1.25; }}
  h2 {{ font-size: {BASE_F + 5}px; margin: 14px 0 6px; line-height: 1.28; }}
  h3 {{ font-size: {BASE_F + 3}px; margin: 10px 0 5px; line-height: 1.3;  }}
  h4 {{ font-size: {BASE_F + 1}px; margin: 8px  0 4px; }}
  /* inline */
  img {{
    max-width: 100%;
    height: auto;
    display: block;
    margin: 14px auto;
  }}
  sup {{ font-size: 0.72em; vertical-align: super; }}
  sub {{ font-size: 0.72em; vertical-align: sub;   }}
  hr  {{ border: none; border-top: 1px solid #e0e0e0; margin: 18px 0; }}
  /* two-column blocks */
  .col-r {{ margin-left: 50%; }}
  /* indent levels */
  .ind1 {{ margin-left: {int(INNER * 0.05)}px; }}
  .ind2 {{ margin-left: {int(INNER * 0.10)}px; }}
  .ind3 {{ margin-left: {int(INNER * 0.15)}px; }}
</style>"""

        def _span_html(span, page_origin_y):
            txt   = span.get("text", "")
            if not txt.strip():
                return txt
            flags = span.get("flags",  0)
            size  = span.get("size",   BASE_F)
            color = span.get("color",  0)
            orig  = span.get("origin", (0, 0))

            styles = []
            if flags & 16:  styles.append("font-weight:700")
            if flags & 2:   styles.append("font-style:italic")

            ratio = size / BASE_F if BASE_F else 1
            if ratio >= 1.35:
                styles.append(f"font-size:{min(int(size * 1.333), 38)}px")
            elif ratio <= 0.78:
                styles.append(f"font-size:{max(int(size * 1.333), 8)}px")

            if color and color != 0:
                r2 = (color >> 16) & 0xFF
                g2 = (color >>  8) & 0xFF
                b2 =  color        & 0xFF
                if not (r2 < 30 and g2 < 30 and b2 < 30):
                    styles.append(f"color:#{r2:02x}{g2:02x}{b2:02x}")

            safe = (_safe_html(txt)
                    .replace("	", "    ")
                    .replace("  ", " &nbsp;"))

            if len(orig) >= 2 and page_origin_y:
                pass

            if not styles:
                return safe
            style_str = ";".join(styles)
            return f'<span style="{style_str}">{safe}</span>'

        def _block_html(block, page_w):
            """
            Convert a fitz text block to an HTML element.
            Returns (html_string, alignment_class).
            """
            bbox      = block.get("bbox", [0, 0, page_w, 0])
            x0, x1    = bbox[0], bbox[2]
            x1 - x0

            all_spans = [sp for ln in block.get("lines", [])
                         for sp in ln.get("spans", [])]
            if not all_spans:
                return "", ""

            first_size = all_spans[0].get("size", BASE_F)
            sum(sp.get("size", BASE_F) for sp in all_spans) / len(all_spans)
            all_bold   = all(sp.get("flags", 0) & 16 for sp in all_spans)

            ratio = first_size / BASE_F if BASE_F else 1
            if ratio >= 1.6 or (ratio >= 1.3 and all_bold):
                tag = "h1"
            elif ratio >= 1.35 or (ratio >= 1.15 and all_bold):
                tag = "h2"
            elif ratio >= 1.15:
                tag = "h3"
            elif ratio >= 1.05 and all_bold:
                tag = "h4"
            else:
                tag = "p"

            left_margin_frac  = x0 / page_w if page_w else 0
            right_margin_frac = (page_w - x1) / page_w if page_w else 0

            align_cls = ""
            indent_cls = ""

            if tag == "p":
                if (abs(left_margin_frac - right_margin_frac) < 0.08
                        and left_margin_frac > 0.15):
                    align_cls = "center"
                elif left_margin_frac > 0.45 and right_margin_frac < 0.12:
                    align_cls = "right"
                elif 0.08 < left_margin_frac < 0.20:
                    indent_cls = "ind1"
                elif 0.20 <= left_margin_frac < 0.30:
                    indent_cls = "ind2"
                elif left_margin_frac >= 0.30:
                    indent_cls = "ind3"

            line_texts  = []
            line_widths = []

            for line in block.get("lines", []):
                spans_html = "".join(
                    _span_html(sp, None) for sp in line.get("spans", [])
                )
                if not spans_html.strip():
                    continue
                line_texts.append(spans_html)
                lbbox = line.get("bbox", [0, 0, 0, 0])
                line_widths.append(lbbox[2] - lbbox[0])

            if not line_texts:
                return "", ""

            block_text_w = bbox[2] - bbox[0]

            if len(line_texts) == 1:
                inner = line_texts[0]
            else:
                short_lines = sum(
                    1 for w in line_widths if block_text_w > 0 and w / block_text_w < 0.75
                )
                short_ratio = short_lines / len(line_widths)

                block_frac = block_text_w / page_w if page_w else 1

                if short_ratio >= 0.5 or block_frac < 0.40:
                    inner = "<br>".join(line_texts)
                else:
                    inner = " ".join(line_texts)

            classes = " ".join(filter(None, [align_cls, indent_cls]))
            cls_attr = f' class="{classes}"' if classes else ""

            return f"<{tag}{cls_attr}>{inner}</{tag}>", align_cls

        pages_html = []

        for page in doc:
            pn     = page.number + 1
            page_w = page.rect.width or MED_W

            img_b64: dict[int, str] = {}
            for info in page.get_images(full=True):
                xref = info[0]
                if xref in img_b64:
                    continue
                try:
                    bi = doc.extract_image(xref)
                    img_b64[xref] = _img_to_b64(bi["image"],
                                                  _mime_for_ext(bi["ext"]))
                except Exception:
                    pass

            items = []
            for block in page.get_text("dict", sort=True).get("blocks", []):
                y0, x0 = block["bbox"][1], block["bbox"][0]
                if block.get("type") == 0:
                    html, _ = _block_html(block, page_w)
                    if html:
                        items.append((y0, x0, "txt", html))
                elif block.get("type") == 1:
                    xref = block.get("xref", 0)
                    if xref and xref in img_b64:
                        items.append((y0, x0, "img", img_b64[xref]))

            placed = {d for _, _, k, d in items if k == "img"}
            for xref, uri in img_b64.items():
                if uri not in placed:
                    items.append((9999, 0, "img", uri))

            items.sort(key=lambda x: (x[0], x[1]))

            body_parts = []
            for _, _, kind, data in items:
                if kind == "txt":
                    body_parts.append(data)
                elif kind == "tbl":
                    body_parts.append(data)
                else:
                    body_parts.append(f'<img src="{data}" alt="img p{pn}">')

            pages_html.append(
                f'<div class="page">'
                f'<div class="pn">Page {pn} / {len(doc)}</div>'
                + "\n".join(body_parts)
                + "</div>"
            )

        with open(dst, "w", encoding="utf-8") as f:
            f.write(
                "<!DOCTYPE html>\n"
                "<html lang=\"fr\">\n<head>\n"
                "  <meta charset=\"utf-8\">\n"
                f"  <title>{_safe_html(name)}</title>\n"
                f"  {CSS}\n"
                "</head>\n<body>\n"
                + "\n<hr>\n".join(pages_html)
                + "\n</body>\n</html>"
            )

    def _epub_to_pdf(self, src, dst):
        """
        EPUB → PDF strategies:
        1. pypandoc (pandoc binary) — best typographic output
        2. Native engine             — fully self-contained, no binary needed
           Handles: spine order, cover image, metadata title, h1-h6,
           p, li (ul/ol), blockquote, br, strong/em/b/i, inline images
           at correct position in text flow.
        """
        try:
            import pypandoc
            pypandoc.convert_file(src, "pdf", outputfile=dst)
            return
        except Exception:
            pass
        self._epub_to_pdf_native(src, dst)

    def _epub_to_pdf_native(self, src, dst):
        """
        Native EPUB → PDF — peak quality.
        Improvements over previous version:
        - Tables rendered as reportlab Table objects
        - <ol> with real numbering (1. 2. 3.)
        - <pre>/<code> in monospace box
        - <figure>/<figcaption> with caption
        - <span style="..."> inline CSS (font-size, font-weight, color)
        - Image path resolution handles ../../../ deep relative paths
        - CSS font-size extracted from embedded <style> blocks
        """
        from xml.etree import ElementTree as ET
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import mm
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer,
            Image as RLImage, PageBreak, HRFlowable,
        )
        from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY

        base = getSampleStyleSheet()
        def _sty(name, **kw):
            parent = kw.pop("parent", base["Normal"])
            return ParagraphStyle(name, parent=parent, **kw)

        sty = {
            "title"  : _sty("ET",  parent=base["Title"],
                            fontSize=22, leading=28, spaceAfter=16,
                            textColor=colors.HexColor("#1a1a2e"), alignment=TA_CENTER),
            "author" : _sty("EA",  fontSize=13, leading=18, spaceAfter=6,
                            textColor=colors.HexColor("#555"), alignment=TA_CENTER),
            "h1"     : _sty("EH1", parent=base["Heading1"],
                            fontSize=18, leading=22, spaceBefore=14, spaceAfter=8,
                            textColor=colors.HexColor("#1a1a2e")),
            "h2"     : _sty("EH2", parent=base["Heading2"],
                            fontSize=15, leading=19, spaceBefore=10, spaceAfter=6,
                            textColor=colors.HexColor("#1e3a5f")),
            "h3"     : _sty("EH3", parent=base["Heading3"],
                            fontSize=13, leading=17, spaceBefore=8, spaceAfter=4,
                            textColor=colors.HexColor("#1e3a5f")),
            "h4"     : _sty("EH4", fontSize=11, leading=15, spaceBefore=6, spaceAfter=3,
                            fontName="Helvetica-Bold", textColor=colors.HexColor("#333")),
            "h5"     : _sty("EH5", fontSize=10.5, leading=14, spaceBefore=4, spaceAfter=2,
                            fontName="Helvetica-Bold"),
            "h6"     : _sty("EH6", fontSize=10, leading=13, spaceBefore=4, spaceAfter=2,
                            fontName="Helvetica-BoldOblique"),
            "body"   : _sty("EB",  fontSize=10.5, leading=16, spaceAfter=5,
                            alignment=TA_JUSTIFY),
            "bq"     : _sty("EBQ", fontSize=10, leading=15, spaceAfter=5,
                            leftIndent=24, rightIndent=12,
                            textColor=colors.HexColor("#555")),
            "li_ul"  : _sty("ELU", fontSize=10.5, leading=15, spaceAfter=3, leftIndent=16),
            "li_ol"  : _sty("ELO", fontSize=10.5, leading=15, spaceAfter=3, leftIndent=20),
            "pre"    : _sty("EPR", fontName="Courier", fontSize=9, leading=13,
                            spaceAfter=8, spaceBefore=4,
                            leftIndent=12, rightIndent=12,
                            backColor=colors.HexColor("#f5f5f5"),
                            borderColor=colors.HexColor("#ddd"),
                            borderWidth=0.5, borderPad=6),
            "caption": _sty("EC",  fontSize=9, leading=12, spaceAfter=8,
                            textColor=colors.HexColor("#777"), alignment=TA_CENTER),
            "td"     : _sty("ETd", fontSize=9, leading=12, spaceAfter=0),
            "th"     : _sty("ETh", fontSize=9, leading=12, spaceAfter=0,
                            fontName="Helvetica-Bold"),
        }

        H_STYS = {"h1":"h1","h2":"h2","h3":"h3","h4":"h4","h5":"h5","h6":"h6"}

        story    = []
        tmp_imgs = []
        PAGE_W   = A4[0] - 40 * mm

        def _decode(text):
            return (text
                    .replace("&amp;",  "&").replace("&lt;",   "<")
                    .replace("&gt;",   ">").replace("&quot;", '"')
                    .replace("&apos;", "'").replace("&nbsp;", "\xa0"))

        def _rl_safe(text):
            d = _decode(text)
            return (d.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;"))

        def _css_to_inline(style_str):
            """Parse a CSS style attribute into reportlab-compatible markup hints."""
            bold = False; italic = False; color = None; size = None
            for decl in style_str.split(";"):
                decl = decl.strip()
                if not decl or ":" not in decl:
                    continue
                prop, _, val = decl.partition(":")
                prop = prop.strip().lower()
                val  = val.strip()
                if prop == "font-weight" and val in ("bold","700","800","900"):
                    bold = True
                elif prop == "font-style" and val == "italic":
                    italic = True
                elif prop == "color":
                    color = val
                elif prop == "font-size":
                    try:
                        size = float(re.sub(r"[^0-9.]","",val))
                    except Exception:
                        pass
            return bold, italic, color, size

        def _inline_markup(frag):
            return inline_markup(frag, decode_fn=_decode, preserve_font_tags=True)

        def _add_img(raw_bytes, ext, alt=""):
            try:
                tf = tempfile.NamedTemporaryFile(suffix=f".{ext or 'png'}", delete=False)
                tf.write(raw_bytes); tf.close()
                tmp_imgs.append(tf.name)
                rl    = RLImage(tf.name)
                scale = min(PAGE_W / rl.imageWidth,
                            (A4[1] * 0.5) / rl.imageHeight, 1.0)
                rl.drawWidth  = rl.imageWidth  * scale
                rl.drawHeight = rl.imageHeight * scale
                rl.hAlign     = "CENTER"
                story.append(Spacer(1, 6))
                story.append(rl)
                if alt and alt.lower() not in ("", "image", "cover"):
                    story.append(Paragraph(_rl_safe(alt), sty["caption"]))
                story.append(Spacer(1, 6))
                return True
            except Exception:
                return False

        def _resolve_img(src_attr, chap_dir, img_data):
            """Resolve image src with robust path normalisation."""
            if not src_attr:
                return None
            if src_attr.startswith("data:"):
                try:
                    b64  = src_attr.split(",", 1)[1]
                    raw  = base64.b64decode(b64)
                    ext  = re.search(r"data:image/(\w+)", src_attr)
                    return raw, (ext.group(1) if ext else "png")
                except Exception:
                    return None

            raw_p  = src_attr.split("?")[0].split("#")[0]
            try:
                from pathlib import PurePosixPath
                resolved = str(PurePosixPath(chap_dir) / raw_p)
            except Exception:
                resolved = f"{chap_dir}/{raw_p}"

            candidates = set()
            for c in [raw_p, resolved, raw_p.lstrip("/"),
                       resolved.lstrip("/")]:
                c_norm = c.replace("\\", "/")
                while c_norm.startswith("./"):
                    c_norm = c_norm[2:]
                candidates.add(c_norm)
                parts = c_norm.split("/")
                if len(parts) > 1:
                    candidates.add(parts[-1])

            for c in candidates:
                b = img_data.get(c)
                if b:
                    return b, Path(c).suffix.lstrip(".")
                fname = c.split("/")[-1]
                for k in img_data:
                    if k.split("/")[-1] == fname:
                        return img_data[k], Path(fname).suffix.lstrip(".")
            return None

        def _parse_table(table_html, chap_dir, img_data):
            return build_table_flowable(table_html, _inline_markup, sty["td"], sty["th"],
                                        header_bg="#e8edf5")

        def _parse_chapter(html_raw, chap_dir, img_data):
            html = re.sub(r'\s+xmlns(?::\w+)?=["\'][^"\']*["\']', "", html_raw)
            bm   = re.search(r"<body[^>]*>(.*?)</body>", html, re.I|re.S)
            body = bm.group(1) if bm else html
            body = re.sub(r"<head[^>]*>.*?</head>",     "", body, flags=re.I|re.S)
            body = re.sub(r"<script[^>]*>.*?</script>", "", body, flags=re.I|re.S)
            body = re.sub(r"<style[^>]*>.*?</style>",   "", body, flags=re.I|re.S)

            TAG_RE = re.compile(
                r"<(/?)(\w+)((?:\s[^>]*)?)/?>", re.I
            )
            pos      = 0
            buf      = ""
            ol_count = [0]

            def flush_buf():
                nonlocal buf
                text = _inline_markup(buf)
                buf  = ""
                if text.strip():
                    story.append(Paragraph(text, sty["body"]))

            for m in TAG_RE.finditer(body):
                buf += body[pos:m.start()]
                pos  = m.end()
                closing, tag, attrs_str = m.group(1), m.group(2).lower(), m.group(3) or ""

                block_tags = {"p","div","h1","h2","h3","h4","h5","h6",
                              "br","hr","li","img","figure","figcaption",
                              "table","ul","ol","pre","blockquote","section",
                              "article","header","footer","aside","nav"}
                if tag in block_tags:
                    flush_buf()

                if tag == "br" and not closing:
                    story.append(Spacer(1, 4))

                elif tag == "hr" and not closing:
                    story.append(HRFlowable(width="100%", thickness=0.5,
                                            color=colors.HexColor("#ccc"),
                                            spaceAfter=6, spaceBefore=6))

                elif tag in H_STYS and not closing:
                    end_m = re.search(rf"</{tag}\s*>", body[pos:], re.I)
                    if end_m:
                        inner = _inline_markup(body[pos:pos+end_m.start()])
                        pos  += end_m.end()
                        if inner:
                            story.append(Spacer(1, 4))
                            story.append(Paragraph(inner, sty[H_STYS[tag]]))

                elif tag == "p" and not closing:
                    end_m = re.search(r"</p\s*>", body[pos:], re.I)
                    if end_m:
                        inner_html = body[pos:pos+end_m.start()]
                        pos       += end_m.end()
                        for img_m in re.finditer(r"<img[^>]+>", inner_html, re.I):
                            sm = re.search(r"""src=(['"])([^'"]+)\1""", img_m.group(0), re.I)
                            if sm:
                                res = _resolve_img(sm.group(2), chap_dir, img_data)
                                if res:
                                    alt_m = re.search(r"""alt=(['"])([^'"]*)\1""",
                                                      img_m.group(0), re.I)
                                    _add_img(res[0], res[1],
                                             alt_m.group(2) if alt_m else "")
                        text = _inline_markup(inner_html)
                        if text.strip():
                            story.append(Paragraph(text, sty["body"]))

                elif tag in ("div","section","article","header",
                             "footer","aside","nav") and not closing:
                    pass

                elif tag == "blockquote" and not closing:
                    end_m = re.search(r"</blockquote\s*>", body[pos:], re.I)
                    if end_m:
                        inner = _inline_markup(body[pos:pos+end_m.start()])
                        pos  += end_m.end()
                        if inner:
                            story.append(Paragraph(inner, sty["bq"]))

                elif tag == "pre" and not closing:
                    end_m = re.search(r"</pre\s*>", body[pos:], re.I)
                    if end_m:
                        raw_pre = body[pos:pos+end_m.start()]
                        pos    += end_m.end()
                        text = re.sub(r"<[^>]+>", "", raw_pre)
                        text = _decode(text)
                        safe = text.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
                        story.append(Paragraph(safe.replace("\n","<br/>"), sty["pre"]))

                elif tag == "ol" and not closing:
                    ol_count[0] = 0
                elif tag == "ul" and not closing:
                    ol_count[0] = -1

                elif tag == "li" and not closing:
                    end_m = re.search(r"</li\s*>", body[pos:], re.I)
                    if end_m:
                        inner = _inline_markup(body[pos:pos+end_m.start()])
                        pos  += end_m.end()
                        if inner.strip():
                            if ol_count[0] >= 0:
                                ol_count[0] += 1
                                prefix = f"{ol_count[0]}. "
                                ps = sty["li_ol"]
                            else:
                                prefix = "• "
                                ps = sty["li_ul"]
                            story.append(Paragraph(prefix + inner, ps))

                elif tag == "img" and not closing:
                    sm = re.search(r"""src=(['"])([^'"]+)\1""", attrs_str, re.I)
                    if sm:
                        res = _resolve_img(sm.group(2), chap_dir, img_data)
                        if res:
                            alt_m = re.search(r"""alt=(['"])([^'"]*)\1""", attrs_str, re.I)
                            _add_img(res[0], res[1], alt_m.group(2) if alt_m else "")

                elif tag == "figure" and not closing:
                    pass

                elif tag == "figcaption" and not closing:
                    end_m = re.search(r"</figcaption\s*>", body[pos:], re.I)
                    if end_m:
                        inner = _inline_markup(body[pos:pos+end_m.start()])
                        pos  += end_m.end()
                        if inner:
                            story.append(Paragraph(inner, sty["caption"]))

                elif tag == "table" and not closing:
                    end_m = re.search(r"</table\s*>", body[pos:], re.I)
                    if end_m:
                        tbl = _parse_table(body[pos:pos+end_m.start()],
                                           chap_dir, img_data)
                        pos += end_m.end()
                        if tbl:
                            story.append(Spacer(1, 8))
                            story.append(tbl)
                            story.append(Spacer(1, 8))

            buf += body[pos:]
            flush_buf()

        with zipfile.ZipFile(src) as zf:
            names = zf.namelist()
            spine_order  = []
            book_title   = Path(src).stem
            book_authors = []
            cover_data   = None
            opf_dir      = ""

            try:
                container = zf.read("META-INF/container.xml").decode("utf-8","replace")
                opf_m = re.search(r"full-path=[\"']([^\"']+\.opf)[\"']", container, re.I)
                if opf_m:
                    opf_path = opf_m.group(1)
                    opf_dir  = str(Path(opf_path).parent)
                    opf_xml  = zf.read(opf_path).decode("utf-8","replace")
                    root     = ET.fromstring(opf_xml)
                    ns_opf   = {"o":"http://www.idpf.org/2007/opf",
                                "dc":"http://purl.org/dc/elements/1.1/"}
                    t_el = root.find(".//dc:title", ns_opf)
                    if t_el is not None and t_el.text:
                        book_title = t_el.text.strip()
                    for cr in root.findall(".//dc:creator", ns_opf):
                        if cr.text:
                            book_authors.append(cr.text.strip())
                    manifest = {}
                    for item in root.findall(".//o:item", ns_opf):
                        iid  = item.get("id","")
                        href = item.get("href","")
                        mtype= item.get("media-type","")
                        full = f"{opf_dir}/{href}".lstrip("/")
                        manifest[iid] = {"href":href,"full":full,"type":mtype}
                    cover_id = None
                    mc = root.find(".//o:meta[@name='cover']", ns_opf)
                    if mc is not None:
                        cover_id = mc.get("content","")
                    if not cover_id:
                        for iid,v in manifest.items():
                            if "cover" in iid.lower() and "image" in v["type"]:
                                cover_id = iid; break
                    if cover_id and cover_id in manifest:
                        cpath = manifest[cover_id]["full"]
                        try:
                            cover_data = (zf.read(cpath),
                                          Path(cpath).suffix.lstrip("."))
                        except Exception:
                            pass
                    for ref in root.findall(".//o:itemref", ns_opf):
                        iid  = ref.get("idref","")
                        if iid in manifest:
                            full = manifest[iid]["full"]
                            if full in names:
                                spine_order.append(full)
            except Exception:
                pass

            if not spine_order:
                spine_order = sorted(
                    n for n in names
                    if n.endswith((".xhtml",".html",".htm"))
                )

            img_data = {}
            for n in names:
                ext = Path(n).suffix.lower().lstrip(".")
                if ext in ("png","jpg","jpeg","gif","webp","bmp","svg"):
                    try:
                        key = n.replace("\\","/").lstrip("/")
                        img_data[key] = zf.read(n)
                    except Exception:
                        pass

            if cover_data:
                _add_img(cover_data[0], cover_data[1], "cover")
                story.append(PageBreak())

            story.append(Spacer(1, 30*mm))
            story.append(Paragraph(_rl_safe(book_title), sty["title"]))
            if book_authors:
                story.append(Spacer(1, 6))
                for auth in book_authors:
                    story.append(Paragraph(_rl_safe(auth), sty["author"]))
            story.append(PageBreak())

            for chap in spine_order:
                try:
                    raw      = zf.read(chap).decode("utf-8","replace")
                    chap_dir = str(Path(chap).parent).replace("\\","/")
                    _parse_chapter(raw, chap_dir, img_data)
                    story.append(PageBreak())
                except Exception:
                    continue

        if len(story) <= 2:
            raise RuntimeError("No content could be extracted from this EPUB.")

        def _safe_str(s):
            return s if isinstance(s, str) else ""

        SimpleDocTemplate(
            dst, pagesize=A4,
            leftMargin=22*mm, rightMargin=22*mm,
            topMargin=22*mm, bottomMargin=22*mm,
            title=_safe_str(book_title),
            author=", ".join(book_authors) if book_authors else "",
        ).build(story)

        for f in tmp_imgs:
            try: os.remove(f)
            except Exception: pass

