"""Document converters: txt, rtf, csv, json, xlsx, pptx."""

from __future__ import annotations

import csv
import json
import re

class DocumentConverters:
    """Document conversion methods for AdvancedConverterEngine."""

    def _txt_to_pdf(self, src: str, dst: str) -> bool:
        """Convert plain text to PDF with word wrapping."""
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.pdfgen import canvas as rl_canvas
            from reportlab.lib.utils import simpleSplit

            c = rl_canvas.Canvas(dst, pagesize=A4)
            width, height = A4
            margin = 40
            y = height - margin
            c.setFont("Helvetica", 11)

            with open(src, "r", encoding="utf-8", errors="replace") as f:
                for line in f:
                    wrapped = simpleSplit(line.rstrip(), "Helvetica", 11, width - 2 * margin)
                    for wl in wrapped:
                        if y < margin:
                            c.showPage()
                            c.setFont("Helvetica", 11)
                            y = height - margin
                        c.drawString(margin, y, wl)
                        y -= 14
            c.save()
            return True
        except Exception as e:
            print(f"[txt→pdf] {e}")
            return False

    def _rtf_to_pdf(self, src: str, dst: str) -> bool:
        """Convert RTF to PDF via internal parser + reportlab."""
        try:
            raw = open(src, "r", encoding="utf-8", errors="replace").read()
            spans = self._rtf_to_spans(raw)
            return self._rtf_spans_to_pdf(spans, dst)
        except Exception as e:
            print(f"[rtf→pdf] {e}")
            return False

    def _rtf_to_spans(self, raw: str) -> list:
        color_table = self._rtf_parse_colortbl(raw)
        tokens = self._rtf_tokenize(raw)
        return self._rtf_parse_spans(tokens, color_table)

    def _rtf_parse_colortbl(self, raw: str) -> list:
        m = re.search(r'\{\\colortbl;([^}]*)\}', raw, re.DOTALL)
        if not m:
            return []
        colors = []
        for part in m.group(1).split(';'):
            nums = re.findall(r'(\d+)', part)
            if len(nums) >= 3:
                colors.append((int(nums[0]), int(nums[1]), int(nums[2])))
        return colors

    def _rtf_tokenize(self, raw: str) -> list:
        tokens = []
        i = 0
        while i < len(raw):
            if raw[i] == '{':
                tokens.append(('{', ''))
                i += 1
            elif raw[i] == '}':
                tokens.append(('}', ''))
                i += 1
            elif raw[i] == '\\':
                m = re.match(r'\\([a-zA-Z]+)(-?\d+)?\s?', raw[i:])
                if m:
                    tokens.append(('cmd', m.group(1), m.group(2)))
                    i += m.end()
                elif raw[i:i+2] == '\\\'':
                    tokens.append(('char', chr(int(raw[i+2:i+4], 16))))
                    i += 4
                else:
                    tokens.append(('char', raw[i+1]))
                    i += 2
            elif raw[i] == '\r\n' or raw[i] == '\n':
                tokens.append(('char', '\n'))
                i += 1
            else:
                tokens.append(('char', raw[i]))
                i += 1
        return tokens

    def _rtf_parse_spans(self, tokens: list, color_table: list) -> list:
        class State:
            def __init__(self):
                self.bold = False
                self.italic = False
                self.font_size = 11
                self.color = None
            def copy(self):
                s = State()
                s.bold = self.bold
                s.italic = self.italic
                s.font_size = self.font_size
                s.color = self.color
                return s

        def _span(text="", par=False, cell_end=False, row_end=False):
            return {"text": text, "bold": state.bold, "italic": state.italic,
                    "size": state.font_size, "par": par, "cell_end": cell_end, "row_end": row_end}

        state = State()
        spans = []
        buf = []
        stack = []

        for tok in tokens:
            if tok[0] == '{':
                stack.append(state)
                state = state.copy()
            elif tok[0] == '}':
                if stack:
                    state = stack.pop()
                if buf:
                    spans.append(_span("".join(buf)))
                    buf = []
            elif tok[0] == 'cmd':
                cmd = tok[1]
                val = tok[2]
                if cmd == 'b':
                    state.bold = val != '0' if val else True
                elif cmd == 'i':
                    state.italic = val != '0' if val else True
                elif cmd == 'fs':
                    if val:
                        state.font_size = max(6, int(val) // 2)
                elif cmd in ('par', 'pard'):
                    if buf:
                        spans.append(_span("".join(buf), par=True))
                        buf = []
                    else:
                        spans.append(_span("", par=True))
                elif cmd in ('cell', 'cellx'):
                    if buf:
                        spans.append(_span("".join(buf), cell_end=True))
                        buf = []
                elif cmd in ('row', 'rowend'):
                    if buf:
                        spans.append(_span("".join(buf), row_end=True))
                        buf = []
            elif tok[0] == 'char':
                if tok[1] != '\n':
                    buf.append(tok[1])

        if buf:
            spans.append(_span("".join(buf)))
        return spans

    def _rtf_spans_to_pdf(self, spans: list, dst: str) -> bool:
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.pdfgen import canvas as rl_canvas
            from reportlab.lib.utils import simpleSplit

            c = rl_canvas.Canvas(dst, pagesize=A4)
            width, height = A4
            margin = 40
            y = height - margin

            for span in spans:
                if span.get("par"):
                    y -= 14
                    continue
                text = span.get("text", "")
                if not text:
                    continue
                size = span.get("size", 11)
                font = "Helvetica-Bold" if span.get("bold") else "Helvetica-Oblique" if span.get("italic") else "Helvetica"
                try:
                    c.setFont(font, size)
                except Exception:
                    c.setFont("Helvetica", size)
                lines = simpleSplit(text, font, size, width - 2 * margin)
                for line in lines:
                    if y < margin:
                        c.showPage()
                        c.setFont(font, size)
                        y = height - margin
                    c.drawString(margin, y, line)
                    y -= size + 2

            c.save()
            return True
        except Exception as e:
            print(f"[rtf→pdf reportlab] {e}")
            return False

    def _txt_to_docx(self, src: str, dst: str) -> bool:
        try:
            from docx import Document
            doc = Document()
            with open(src, "r", encoding="utf-8", errors="replace") as f:
                for line in f:
                    doc.add_paragraph(line.rstrip())
            doc.save(dst)
            return True
        except Exception as e:
            print(f"[txt→docx] {e}")
            return False

    def _rtf_to_docx(self, src: str, dst: str) -> bool:
        try:
            from docx import Document
            raw = open(src, "r", encoding="utf-8", errors="replace").read()
            spans = self._rtf_to_spans(raw)
            doc = Document()
            para = doc.add_paragraph()
            for span in spans:
                if span.get("par"):
                    para = doc.add_paragraph()
                    continue
                text = span.get("text", "")
                if not text:
                    continue
                run = para.add_run(text)
                run.bold = span.get("bold", False)
                run.italic = span.get("italic", False)
            doc.save(dst)
            return True
        except Exception as e:
            print(f"[rtf→docx] {e}")
            return False

    def _csv_to_json(self, src: str, dst: str) -> bool:
        try:
            with open(src, "r", encoding="utf-8", errors="replace", newline="") as f:
                reader = csv.DictReader(f)
                data = list(reader)
            with open(dst, "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            return True
        except Exception as e:
            print(f"[csv→json] {e}")
            return False

    def _json_to_csv(self, src: str, dst: str) -> bool:
        try:
            with open(src, "r", encoding="utf-8") as f:
                data = json.load(f)
            if not isinstance(data, list) or not data:
                return False
            keys = list(data[0].keys())
            with open(dst, "w", encoding="utf-8", newline="") as f:
                writer = csv.DictWriter(f, fieldnames=keys)
                writer.writeheader()
                writer.writerows(data)
            return True
        except Exception as e:
            print(f"[json→csv] {e}")
            return False

    def _xlsx_to_pdf(self, src: str, dst: str) -> bool:
        try:
            import openpyxl
            from reportlab.lib.pagesizes import A4, landscape
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
            from reportlab.lib import colors

            wb = openpyxl.load_workbook(src, data_only=True)
            elements = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    rows.append([str(c) if c is not None else "" for c in row])
                if rows:
                    t = Table(rows)
                    t.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('GRID', (0, 0), (-1, -1), 1, colors.black),
                        ('FONTSIZE', (0, 0), (-1, -1), 8),
                    ]))
                    elements.append(t)

            if elements:
                doc = SimpleDocTemplate(dst, pagesize=landscape(A4))
                doc.build(elements)
            wb.close()
            return True
        except Exception as e:
            print(f"[xlsx→pdf] {e}")
            return False

    def _xlsx_to_json(self, src: str, dst: str) -> bool:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(src, data_only=True)
            result = {}
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    rows.append([str(c) if c is not None else "" for c in row])
                result[sheet_name] = rows
            wb.close()
            with open(dst, "w", encoding="utf-8") as f:
                json.dump(result, f, ensure_ascii=False, indent=2)
            return True
        except Exception as e:
            print(f"[xlsx→json] {e}")
            return False

    def _xlsx_to_csv(self, src: str, dst: str) -> bool:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(src, data_only=True)
            ws = wb.active
            with open(dst, "w", encoding="utf-8", newline="") as f:
                writer = csv.writer(f)
                for row in ws.iter_rows(values_only=True):
                    writer.writerow([str(c) if c is not None else "" for c in row])
            wb.close()
            return True
        except Exception as e:
            print(f"[xlsx→csv] {e}")
            return False

    def _pptx_to_pdf(self, src: str, dst: str) -> bool:
        try:
            from pptx import Presentation
            from reportlab.lib.pagesizes import landscape, A4
            from reportlab.pdfgen import canvas as rl_canvas

            prs = Presentation(src)
            c = rl_canvas.Canvas(dst, pagesize=landscape(A4))
            width, height = landscape(A4)

            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                c.setFont("Helvetica", 12)
                                c.drawString(50, height - 100, text[:100])
                c.showPage()
            c.save()
            return True
        except Exception as e:
            print(f"[pptx→pdf] {e}")
            return False
