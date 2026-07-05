"""OptimizationMixin — Office file optimization methods."""

import os
from pathlib import Path
from datetime import datetime

class OptimizationMixin:
    """Mixin: office file optimization methods for FileConverterApp."""

    def optimize_office_files(self, office_files, optimization_type, quality_level, remove_metadata, compress_images, keep_backup):
        if hasattr(self, 'active_templates') and 'office_optimization' in self.active_templates:
            del self.active_templates['office_optimization']
        """Optimize office and image files"""
        output_dir = self.get_output_directory()
        if not output_dir:
            return

        IMAGE_EXTS  = {'.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.webp', '.gif'}
        EXCEL_EXTS  = {'.xlsx', '.xls'}
        AUDIO_EXTS  = {'.mp3', '.wav', '.aac', '.flac', '.ogg'}
        VIDEO_EXTS  = {'.mp4', '.avi', '.mkv', '.mov', '.webm'}
        WEB_EXTS    = {'.json', '.html', '.htm'}
        EPUB_EXTS   = {'.epub'}

        self.show_progress(True, self.translate_text("len_off").format(len(office_files)))

        success_count = 0
        total_original_size = 0
        total_compressed_size = 0
        start_time = datetime.now()

        for i, file_path in enumerate(office_files):
            try:
                file_ext = Path(file_path).suffix.lower()
                original_size = os.path.getsize(file_path)
                total_original_size += original_size

                if keep_backup:
                    output_file = os.path.join(output_dir, f"optimized_{Path(file_path).name}")
                else:
                    output_file = file_path

                operation_start = datetime.now()

                if quality_level == 0:
                    compression_level = "high"
                elif quality_level == 1:
                    compression_level = "normal"
                elif quality_level == 2:
                    compression_level = "very_reduced"

                if file_ext == '.pdf':
                    success = self.optimize_pdf_file(file_path, output_file, compression_level, remove_metadata, compress_images)
                elif file_ext in ['.docx', '.doc']:
                    success = self.optimize_word_file(file_path, output_file, compression_level, remove_metadata, compress_images)
                elif file_ext in ['.pptx', '.ppt']:
                    success = self.optimize_powerpoint_file(file_path, output_file, compression_level, remove_metadata, compress_images)
                elif file_ext in EXCEL_EXTS:
                    success = self.optimize_excel_file(file_path, output_file, compression_level, remove_metadata)
                elif file_ext in IMAGE_EXTS:
                    success = self.optimize_image_file(file_path, output_file, quality_level)
                elif file_ext in AUDIO_EXTS:
                    success = self.optimize_av_file(file_path, output_file, quality_level, 'audio')
                elif file_ext in VIDEO_EXTS:
                    success = self.optimize_av_file(file_path, output_file, quality_level, 'video')
                elif file_ext in WEB_EXTS:
                    success = self.optimize_web_file(file_path, output_file, file_ext)
                elif file_ext in EPUB_EXTS:
                    success = self.optimize_epub_file(file_path, output_file, compress_images, quality_level)
                else:
                    success = False

                if success:
                    compressed_size = os.path.getsize(output_file) if os.path.exists(output_file) else original_size
                    total_compressed_size += compressed_size
                    operation_time = (datetime.now() - operation_start).total_seconds()
                    self.db_manager.add_conversion_record(
                        source_file=file_path, source_format=file_ext.upper().replace('.', ''),
                        target_file=output_file, target_format=file_ext.upper().replace('.', ''),
                        operation_type="office_optimization", file_size=original_size,
                        conversion_time=operation_time, success=True,
                        notes=f"Type: {optimization_type}, Quality: {quality_level}")
                    self.achievement_system.record_conversion("office_optimization", original_size, True)
                    success_count += 1
                self.progress_bar.setValue(int((i + 1) / len(office_files) * 100))
            except Exception as e:
                self.db_manager.add_conversion_record(
                    source_file=file_path, source_format=Path(file_path).suffix.upper().replace('.', ''),
                    target_file="", target_format="", operation_type="office_optimization",
                    file_size=0, conversion_time=0, success=False, notes=f"Error: {str(e)}")
                print(f"Optimization error {file_path}: {e}")

        total_time = (datetime.now() - start_time).total_seconds()
        self.show_progress(False)

        if total_original_size > 0:
            compression_rate = ((total_original_size - total_compressed_size) / total_original_size * 100)
            savings_mb = (total_original_size - total_compressed_size) / (1024 * 1024)
            message = self.translate_text("msg_1").format(success_count, len(office_files), f"{total_time:.1f}")
            message += self.translate_text("msg_2").format(f"{savings_mb:.2f}", f"{compression_rate:.1f}")
            message += self.translate_text("msg_3").format(output_dir)
        else:
            message = self.translate_text("msg_4").format(success_count, len(office_files))
        from PySide6.QtWidgets import QMessageBox
        QMessageBox.information(self, self.translate_text("Succès"), self.translate_text(message))

    def optimize_pdf_file(self, pdf_path, output_path, compression_level, remove_metadata, compress_images):
        try:
            import fitz
            pdf_document = fitz.open(pdf_path)
            if compression_level == "very_reduced":
                save_options = {'garbage': 3, 'deflate': True, 'clean': True, 'deflate_images': compress_images, 'deflate_fonts': True}
            elif compression_level == "high":
                save_options = {'garbage': 3, 'deflate': True, 'clean': True, 'deflate_images': compress_images, 'deflate_fonts': True}
            else:
                save_options = {'garbage': 2, 'deflate': True, 'clean': True}
            if remove_metadata:
                pdf_document.set_metadata({})
            pdf_document.save(output_path, **save_options)
            pdf_document.close()
            return True
        except Exception as e:
            print(f"PDF optimization error {pdf_path}: {e}")
            return False

    def optimize_word_file(self, word_path, output_path, compression_level, remove_metadata, compress_images):
        try:
            from docx import Document
            import io as _io
            doc = Document(word_path)
            if compress_images:
                quality_map = {"high": 85, "normal": 75, "very_reduced": 55}
                jpeg_quality = quality_map.get(compression_level, 75)
                try:
                    from PIL import Image as PILImage
                    for rel in doc.part.rels.values():
                        if "image" in rel.reltype:
                            try:
                                blob = rel.target_part.blob
                                ext  = Path(rel.target_part.partname).suffix.lower()
                                if ext in ('.jpg', '.jpeg', '.png', '.bmp', '.tiff'):
                                    buf_in  = _io.BytesIO(blob)
                                    buf_out = _io.BytesIO()
                                    with PILImage.open(buf_in) as img:
                                        img.convert('RGB').save(buf_out, format='JPEG', quality=jpeg_quality, optimize=True)
                                    new_blob = buf_out.getvalue()
                                    if len(new_blob) < len(blob):
                                        rel.target_part._blob = new_blob
                            except Exception:
                                pass
                except ImportError:
                    pass
            if remove_metadata:
                doc.core_properties.title = ""
                doc.core_properties.author = ""
                doc.core_properties.subject = ""
                doc.core_properties.keywords = ""
                doc.core_properties.comments = ""
                doc.core_properties.last_modified_by = ""
            doc.save(output_path)
            return True
        except Exception as e:
            print(f"Word optimization error {word_path}: {e}")
            return False

    def optimize_powerpoint_file(self, ppt_path, output_path, compression_level, remove_metadata, compress_images):
        try:
            from pptx import Presentation
            import io
            prs = Presentation(ppt_path)
            slides_to_remove = [i for i, slide in enumerate(prs.slides) if not slide.shapes]
            xml_slides = prs.slides._sldIdLst
            for i in reversed(slides_to_remove):
                xml_slides.remove(xml_slides[i])
            if compress_images:
                quality_map = {"high": 85, "normal": 75, "very_reduced": 55}
                jpeg_quality = quality_map.get(compression_level, 75)
                try:
                    from PIL import Image as PILImage
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.shape_type == 13:
                                try:
                                    img_part = shape.image
                                    blob = img_part.blob
                                    ext  = img_part.ext.lower()
                                    if ext in ("jpg", "jpeg", "png", "bmp", "tiff"):
                                        buf_in  = io.BytesIO(blob)
                                        buf_out = io.BytesIO()
                                        with PILImage.open(buf_in) as img:
                                            rgb = img.convert("RGB")
                                            rgb.save(buf_out, format="JPEG", quality=jpeg_quality, optimize=True)
                                        new_blob = buf_out.getvalue()
                                        if len(new_blob) < len(blob):
                                            img_part._blob = new_blob
                                except Exception:
                                    pass
                except ImportError:
                    pass
            if remove_metadata:
                prs.core_properties.title = ""
                prs.core_properties.author = ""
                prs.core_properties.subject = ""
                prs.core_properties.keywords = ""
                prs.core_properties.comments = ""
                prs.core_properties.last_modified_by = ""
            prs.save(output_path)
            return True
        except Exception as e:
            print(f"PowerPoint optimization error {ppt_path}: {e}")
            return False

    def optimize_av_file(self, src_path, output_path, quality_level, media_type):
        try:
            import shutil, sys, subprocess
            ffmpeg_bin = shutil.which("ffmpeg")
            if not ffmpeg_bin:
                candidates = [
                    r"C:\ffmpeg\bin\ffmpeg.exe", r"C:\Program Files\ffmpeg\bin\ffmpeg.exe",
                    r"C:\Program Files (x86)\ffmpeg\bin\ffmpeg.exe",
                    os.path.join(os.environ.get("LOCALAPPDATA", ""), "ffmpeg", "bin", "ffmpeg.exe"),
                    os.path.join(os.environ.get("APPDATA", ""), "ffmpeg", "bin", "ffmpeg.exe"),
                    os.path.join(getattr(sys, "_MEIPASS", ""), "ffmpeg.exe"),
                    "/usr/bin/ffmpeg", "/usr/local/bin/ffmpeg",
                ]
                for c in candidates:
                    if c and os.path.isfile(c):
                        ffmpeg_bin = c
                        break
            if not ffmpeg_bin:
                return False
            ext = Path(output_path).suffix.lower().lstrip(".")
            if media_type == 'audio':
                AUDIO_PRESETS = {
                    "mp3": {0: ["-codec:a", "libmp3lame", "-q:a", "2", "-ar", "44100"], 1: ["-codec:a", "libmp3lame", "-q:a", "4", "-ar", "44100"], 2: ["-codec:a", "libmp3lame", "-q:a", "7", "-ar", "44100"]},
                    "aac": {0: ["-codec:a", "aac", "-b:a", "192k", "-ar", "44100"], 1: ["-codec:a", "aac", "-b:a", "128k", "-ar", "44100"], 2: ["-codec:a", "aac", "-b:a", "96k", "-ar", "44100"]},
                    "ogg": {0: ["-codec:a", "libvorbis", "-q:a", "6"], 1: ["-codec:a", "libvorbis", "-q:a", "4"], 2: ["-codec:a", "libvorbis", "-q:a", "2"]},
                    "flac": {0: ["-codec:a", "flac", "-compression_level", "5"], 1: ["-codec:a", "flac", "-compression_level", "8"], 2: ["-codec:a", "flac", "-compression_level", "12"]},
                    "wav": {0: ["-codec:a", "pcm_s16le", "-ar", "44100"], 1: ["-codec:a", "pcm_s16le", "-ar", "44100"], 2: ["-codec:a", "pcm_s16le", "-ar", "22050"]},
                }
                default_audio = {0: ["-codec:a", "libmp3lame", "-q:a", "2"], 1: ["-codec:a", "libmp3lame", "-q:a", "4"], 2: ["-codec:a", "libmp3lame", "-q:a", "7"]}
                presets = AUDIO_PRESETS.get(ext, default_audio)
                args = presets.get(quality_level, presets[1])
            else:
                VIDEO_PRESETS = {
                    "mp4": {0: ["-codec:v", "libx264", "-crf", "18", "-preset", "slow", "-codec:a", "aac", "-b:a", "192k", "-movflags", "+faststart"], 1: ["-codec:v", "libx264", "-crf", "23", "-preset", "medium", "-codec:a", "aac", "-b:a", "128k", "-movflags", "+faststart"], 2: ["-codec:v", "libx264", "-crf", "28", "-preset", "fast", "-codec:a", "aac", "-b:a", "96k", "-movflags", "+faststart"]},
                    "mkv": {0: ["-codec:v", "libx264", "-crf", "18", "-preset", "slow", "-codec:a", "aac", "-b:a", "192k"], 1: ["-codec:v", "libx264", "-crf", "23", "-preset", "medium", "-codec:a", "aac", "-b:a", "128k"], 2: ["-codec:v", "libx264", "-crf", "28", "-preset", "fast", "-codec:a", "aac", "-b:a", "96k"]},
                    "webm": {0: ["-codec:v", "libvpx-vp9", "-crf", "24", "-b:v", "0", "-codec:a", "libopus", "-b:a", "160k"], 1: ["-codec:v", "libvpx-vp9", "-crf", "33", "-b:v", "0", "-codec:a", "libopus", "-b:a", "128k"], 2: ["-codec:v", "libvpx-vp9", "-crf", "42", "-b:v", "0", "-codec:a", "libopus", "-b:a", "96k"]},
                }
                default_video = {0: ["-codec:v", "libx264", "-crf", "18", "-preset", "slow", "-codec:a", "aac", "-b:a", "192k"], 1: ["-codec:v", "libx264", "-crf", "23", "-preset", "medium", "-codec:a", "aac", "-b:a", "128k"], 2: ["-codec:v", "libx264", "-crf", "28", "-preset", "fast", "-codec:a", "aac", "-b:a", "96k"]}
                presets = VIDEO_PRESETS.get(ext, default_video)
                args = presets.get(quality_level, presets[1])
            cmd = [ffmpeg_bin, "-y", "-i", src_path] + args + [output_path]
            _no_window = subprocess.CREATE_NO_WINDOW if hasattr(subprocess, "CREATE_NO_WINDOW") else 0
            result = subprocess.run(cmd, capture_output=True, timeout=3600, creationflags=_no_window)
            if result.returncode != 0:
                return False
            if os.path.exists(output_path):
                orig_size = os.path.getsize(src_path)
                new_size  = os.path.getsize(output_path)
                if new_size >= orig_size * 1.05:
                    shutil.copy2(src_path, output_path)
            return True
        except Exception as e:
            print(f"[optimize_av] Error {src_path}: {e}")
            return False

    def optimize_web_file(self, src_path, output_path, file_ext):
        try:
            with open(src_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read()
            if file_ext == '.json':
                import json as _json
                data = _json.loads(content)
                minified = _json.dumps(data, ensure_ascii=False, separators=(',', ':'))
            else:
                import re as _re
                minified = _re.sub(r'<!--.*?-->', '', content, flags=_re.DOTALL)
                minified = _re.sub(r'>\s+<', '><', minified)
                minified = _re.sub(r'[ \t]{2,}', ' ', minified)
                lines = [l.strip() for l in minified.splitlines()]
                minified = '\n'.join(l for l in lines if l)
            orig_bytes = content.encode('utf-8')
            new_bytes  = minified.encode('utf-8')
            if len(new_bytes) < len(orig_bytes):
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(minified)
            else:
                import shutil as _sh
                _sh.copy2(src_path, output_path)
            return True
        except Exception as e:
            print(f"Web file optimization error {src_path}: {e}")
            return False

    def optimize_epub_file(self, src_path, output_path, compress_images, quality_level):
        try:
            import zipfile, io as _io
            quality_map = {0: 85, 1: 75, 2: 55}
            img_quality = quality_map.get(quality_level, 75)
            IMAGE_EXTS = {'.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp'}
            buf_out = _io.BytesIO()
            with zipfile.ZipFile(src_path, 'r') as zin, zipfile.ZipFile(buf_out, 'w', compression=zipfile.ZIP_DEFLATED, compresslevel=9) as zout:
                for item in zin.infolist():
                    data = zin.read(item.filename)
                    ext  = Path(item.filename).suffix.lower()
                    if compress_images and ext in IMAGE_EXTS and len(data) > 2048:
                        try:
                            from PIL import Image as PILImage
                            buf_img = _io.BytesIO(data)
                            buf_new = _io.BytesIO()
                            with PILImage.open(buf_img) as img:
                                if ext in ('.jpg', '.jpeg'):
                                    img.convert('RGB').save(buf_new, format='JPEG', quality=img_quality, optimize=True)
                                elif ext == '.png':
                                    compress_lvl = max(1, min(9, int(quality_level * 3) + 1))
                                    img.save(buf_new, format='PNG', optimize=True, compress_level=compress_lvl)
                                else:
                                    img.convert('RGB').save(buf_new, format='JPEG', quality=img_quality, optimize=True)
                            new_data = buf_new.getvalue()
                            if len(new_data) < len(data):
                                data = new_data
                        except Exception:
                            pass
                    if item.filename == 'mimetype':
                        zout.writestr(item, data, compress_type=zipfile.ZIP_STORED)
                    else:
                        zout.writestr(item, data)
            new_bytes  = buf_out.getvalue()
            orig_bytes = os.path.getsize(src_path)
            if len(new_bytes) < orig_bytes:
                with open(output_path, 'wb') as f:
                    f.write(new_bytes)
            else:
                import shutil as _sh
                _sh.copy2(src_path, output_path)
            return True
        except Exception as e:
            print(f"EPUB optimization error {src_path}: {e}")
            return False

    def optimize_excel_file(self, xlsx_path, output_path, compression_level, remove_metadata):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(xlsx_path)
            if remove_metadata:
                wb.properties.title = ""
                wb.properties.creator = ""
                wb.properties.subject = ""
                wb.properties.keywords = ""
                wb.properties.description = ""
                wb.properties.lastModifiedBy = ""
            empty_sheets = [ws.title for ws in wb.worksheets if ws.max_row == 1 and ws.max_column == 1 and ws.cell(1, 1).value is None]
            for name in empty_sheets:
                if len(wb.sheetnames) > 1:
                    del wb[name]
            wb.save(output_path)
            return True
        except Exception as e:
            print(f"Excel optimization error {xlsx_path}: {e}")
            return False

    def optimize_image_file(self, img_path, output_path, quality_level):
        try:
            from PIL import Image as PILImage
            quality_map = {0: 85, 1: 75, 2: 55}
            quality = quality_map.get(quality_level, 75)
            ext_map = {".jpg": "JPEG", ".jpeg": "JPEG", ".png": "PNG", ".bmp": "PNG", ".tiff": "TIFF", ".webp": "WEBP", ".gif": "GIF"}
            src_ext  = Path(img_path).suffix.lower()
            dst_ext  = Path(output_path).suffix.lower()
            fmt_out  = ext_map.get(dst_ext, ext_map.get(src_ext, "JPEG"))
            with PILImage.open(img_path) as img:
                exif_data = None
                try:
                    exif_data = img.info.get("exif")
                except Exception:
                    pass
                if fmt_out == "JPEG":
                    rgb = img.convert("RGB")
                    save_kwargs = {"quality": quality, "optimize": True}
                    if exif_data:
                        save_kwargs["exif"] = exif_data
                    rgb.save(output_path, format="JPEG", **save_kwargs)
                elif fmt_out == "PNG":
                    compress = max(1, min(9, int(quality_level * 3)))
                    img.save(output_path, format="PNG", optimize=True, compress_level=compress)
                elif fmt_out == "WEBP":
                    img.save(output_path, format="WEBP", quality=quality, method=6)
                else:
                    img.save(output_path, format=fmt_out)
            return True
        except Exception as e:
            print(f"Image optimization error {img_path}: {e}")
            return False
